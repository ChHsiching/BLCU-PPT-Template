"""QA gate for a rendered .pptx deck (G2 machine checks, pptx side).

Usage:
    python scripts/qa_check_pptx.py <out.pptx> --deck <deck.json>
        [--templates-dir DIR] [--manifest PATH] [--com-screenshots DIR]
        [--format text|json]

The pptx is checked against the deck it claims to render (deck.json stays the
single source of truth). Checks, each a finding on failure:

  count     slide count == deck page count
  titles    per-page title alignment (title placeholder; agenda pages carry
            none, so their title block is matched against the rebuilt
            AgendaLabel box)
  residue   placeholder-text grep over every text run of every slide
            (TODO/FIXME/TBD/xxx/lorem/[insert/placeholder/待补充/…)
  budget    the deck itself revalidated against the manifest
            (scripts/validate_deck.py pure function) — catches decks edited
            after rendering or rendered through the API without validation
  structure python-pptx level: every shape inside the slide canvas;
            text-bearing shapes (pictures, the page-number placeholder and
            in-slot caption strips excepted) above safe_canvas
            .content_bottom_y; zero speaker notes (演讲稿 is a separate
            deliverable); pages referencing images carry at least that many
            embedded pictures (formula fallbacks may add more)
  style     run-level typography.tokens: every content run carries the token
            family on ea+latin (font consistency page to page), role-bound
            shapes match their role's size/weight/color (weight hierarchy;
            **emphasis** runs green-bold only inside body-flow shapes), math
            runs stay Cambria Math at the formula size, and any run color
            outside the token role colors is flagged (band color never a
            text color; accent only as bold)
  brand     every slide inherits a master carrying its manifest brand_layer
            group: the band rect and each logo picture at the measured
            geometry, logo media matching the declared file name

--com-screenshots DIR adds the optional PowerPoint COM spot check: every slide
exported as a PNG into DIR for visual overflow inspection. A missing COM
environment (not Windows, no PowerPoint, no pywin32) is never itself a
finding — a note goes to stderr and the programmatic checks above stand alone.

Exit codes: 0 all green, 1 findings, 2 usage / IO error.
The core is the pure check_pptx(prs, deck, manifest, image_root) -> list
of validate_deck.Finding; scripts/qa_check_web.py reuses it for the export
chain.
"""
from __future__ import annotations

import argparse
import json
import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu

sys.path.insert(0, str(Path(__file__).resolve().parent))
import validate_deck as vd  # noqa: E402

NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
NS_M = "http://schemas.openxmlformats.org/officeDocument/2006/math"
NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"
NS_R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
NS_SVG = "http://schemas.microsoft.com/office/drawing/2016/SVG/main"

PICTURE_SHAPE_TYPE = 13  # MSO_SHAPE_TYPE.PICTURE, avoided as a hard import

# placeholder-residue patterns: (name, matcher). Regexes are for scripts and
# latin filler with word boundaries; CJK fillers match as substrings (CJK has
# no word boundaries to lean on). The list is deliberately conservative:
# bare "占位" / "待定" are excluded because they live in legitimate academic
# terms (可学习占位符 / 待定系数法) — only unambiguous filler compounds match.
_RESIDUE_REGEXES = [
    ("todo", re.compile(r"\btodo\b", re.IGNORECASE)),
    ("fixme", re.compile(r"\bfixme\b", re.IGNORECASE)),
    ("tbd", re.compile(r"\btbd\b", re.IGNORECASE)),
    ("xxx", re.compile(r"x{3,}", re.IGNORECASE)),
    ("lorem", re.compile(r"lorem", re.IGNORECASE)),
    ("insert-tag", re.compile(r"[\[<\{]insert\b", re.IGNORECASE)),
    ("placeholder", re.compile(r"\bplaceholder\b", re.IGNORECASE)),
    ("sample-text", re.compile(r"\bsample text\b", re.IGNORECASE)),
]
_RESIDUE_SUBSTRINGS = ("待补充", "待填写", "此处填写", "示例文本")

SNIPPET_RADIUS = 10


def find_residue(text: str) -> list[tuple[str, str]]:
    """Return [(pattern_name, snippet)] for every filler marker in text."""
    hits = []
    for name, rx in _RESIDUE_REGEXES:
        for m in rx.finditer(text):
            lo, hi = max(0, m.start() - SNIPPET_RADIUS), min(len(text), m.end() + SNIPPET_RADIUS)
            hits.append((name, text[lo:hi]))
    for needle in _RESIDUE_SUBSTRINGS:
        start = text.find(needle)
        if start != -1:
            lo, hi = max(0, start - SNIPPET_RADIUS), min(len(text), start + len(needle) + SNIPPET_RADIUS)
            hits.append((needle, text[lo:hi]))
    return hits


# ---------------------------------------------------------------------------
# slide introspection (mirrors the renderer's own serialization)
# ---------------------------------------------------------------------------


def ph_type(sp) -> str | None:
    """Placeholder type of a shape, None when not a placeholder."""
    ph = sp._element.find(f".//{{{NS_P}}}ph")
    if ph is None:
        return None
    return ph.get("type") or "body"


def shape_text(sp) -> str:
    """All visible text of a shape: a:t and m:t alike, concatenated."""
    parts = []
    for t in sp._element.iter(f"{{{NS_A}}}t"):
        parts.append(t.text or "")
    for t in sp._element.iter(f"{{{NS_M}}}t"):
        parts.append(t.text or "")
    return "".join(parts)


def slide_title_text(slide) -> str | None:
    """Title text of a rendered slide.

    Content pages carry a title placeholder; agenda pages have none — their
    title block lands in the renderer's rebuilt AgendaLabel box.
    """
    for sp in slide.shapes:
        if ph_type(sp) == "title":
            return sp.text_frame.text
    for sp in slide.shapes:
        if sp.name == "AgendaLabel":
            return sp.text_frame.text
    return None


def _deck_title(page) -> str | None:
    """Title text of a deck page, tolerating schema-broken pages (None).

    The tool's job is to judge exactly the decks that may have been edited
    after rendering, so traversal never raises on missing fields — the
    validator's schema findings report those, this just skips them.
    """
    if not isinstance(page, dict):
        return None
    blocks = page.get("blocks")
    if not isinstance(blocks, list):
        return None
    for block in blocks:
        if isinstance(block, dict) and block.get("type") == "title" \
                and isinstance(block.get("text"), str):
            return block["text"]
    return None


def _deck_image_count(page) -> int:
    if not isinstance(page, dict):
        return 0
    blocks = page.get("blocks")
    if not isinstance(blocks, list):
        return 0
    return sum(1 for b in blocks if isinstance(b, dict) and b.get("type") == "image")


# ---------------------------------------------------------------------------
# core checks (pure; reused by qa_check_web's export chain)
# ---------------------------------------------------------------------------


def check_pptx(prs, deck: dict, manifest: dict, image_root: Path) -> list:
    """Run every pptx-side check; return findings in slide order."""
    findings: list[vd.Finding] = []

    # budget: the deck is the single source of truth — revalidate it
    findings.extend(vd.validate_deck(deck, manifest, image_root))

    # count (deck pages may be anything; _deck_* helpers stay total)
    slides = list(prs.slides)
    pages = deck.get("pages") if isinstance(deck, dict) else None
    pages = pages if isinstance(pages, list) else []

    # count
    if len(slides) != len(pages):
        findings.append(vd.Finding(
            "", "pptx.page_count",
            f"pptx has {len(slides)} slide(s) but the deck declares {len(pages)} page(s)",
        ))

    # per-slide checks over the aligned prefix
    for i, slide in enumerate(slides):
        spath = f"slides[{i + 1}]"
        page = pages[i] if i < len(pages) else None

        # titles (skip pages whose deck side has no usable title block —
        # the validator's schema.missing_title finding already covers them)
        if page is not None:
            expected = _deck_title(page)
            actual = slide_title_text(slide)
            if expected is not None:
                if actual is None:
                    findings.append(vd.Finding(
                        spath, "pptx.title_missing",
                        f"no title found on slide {i + 1} (expected: {expected!r})",
                    ))
                elif expected != actual:
                    findings.append(vd.Finding(
                        spath, "pptx.title_mismatch",
                        f"title on slide {i + 1} is {actual!r} but the deck declares {expected!r}",
                    ))

        # residue + structure, shape by shape
        n_pictures = 0
        for sp in slide.shapes:
            is_picture = sp.shape_type == PICTURE_SHAPE_TYPE
            if is_picture:
                n_pictures += 1
            text = shape_text(sp)
            if text:
                for name, snippet in find_residue(text):
                    findings.append(vd.Finding(
                        f"{spath}.{sp.name}", "pptx.residue",
                        f"placeholder residue [{name}] on slide {i + 1}: …{snippet}…",
                    ))
            findings.extend(_check_geometry(sp, spath, manifest, is_picture))

        # style + brand (#20): run-level tokens and the inherited master
        findings.extend(_check_style(slide, page, spath, manifest))
        findings.extend(_check_brand(slide, page, spath, manifest))

        if page is not None:
            n_images = _deck_image_count(page)
            if n_pictures < n_images:
                findings.append(vd.Finding(
                    spath, "pptx.images_missing",
                    f"slide {i + 1} carries {n_pictures} picture(s) but the deck "
                    f"references {n_images} image block(s)",
                ))

        if slide.has_notes_slide:
            findings.append(vd.Finding(
                spath, "pptx.notes_slide",
                f"slide {i + 1} has speaker notes; 演讲稿 is a separate deliverable "
                f"and never lives in the pptx",
            ))

    return findings


def _check_geometry(sp, spath: str, manifest: dict, is_picture: bool) -> list:
    """Canvas bounds for every shape; safe content bottom for text shapes.

    Exceptions mirror manifest.safe_canvas: pictures only answer to the canvas,
    the page-number placeholder sits in its own reserved band, and caption
    strips overlay the internal bottom edge of their owning image slot.
    """
    findings: list[vd.Finding] = []
    if sp.left is None or sp.top is None or sp.width is None or sp.height is None:
        return findings
    tol = 0.02  # absorbs int-EMU rounding, nothing more
    canvas_w, canvas_h = manifest["slide_size"]["w"], manifest["slide_size"]["h"]
    left, top = Emu(sp.left).inches, Emu(sp.top).inches
    right = Emu(sp.left + sp.width).inches
    bottom = Emu(sp.top + sp.height).inches
    if left < -tol or top < -tol or right > canvas_w + tol or bottom > canvas_h + tol:
        findings.append(vd.Finding(
            f"{spath}.{sp.name}", "pptx.canvas_overflow",
            f"'{sp.name}' spans ({left:.2f}, {top:.2f})-({right:.2f}, {bottom:.2f}); "
            f"outside the {canvas_w}x{canvas_h} canvas",
        ))
        return findings

    if is_picture or ph_type(sp) == "sldNum" or sp.name.startswith("Caption"):
        return findings
    content_bottom = manifest.get("safe_canvas", {}).get("content_bottom_y")
    if content_bottom is not None and bottom > content_bottom + tol:
        findings.append(vd.Finding(
            f"{spath}.{sp.name}", "pptx.content_bottom",
            f"'{sp.name}' bottom edge {bottom:.2f} crosses safe_canvas "
            f"content_bottom_y {content_bottom}",
        ))
    return findings


# ---------------------------------------------------------------------------
# style layer (#20): typography.tokens at run level + master brand layer
# ---------------------------------------------------------------------------

# renderer textbox name -> (role_bindings region key, emphasis runs allowed).
# Both columns mirror what render_pptx actually consumes: fill_text_formula
# styles TextArea and TextFullArea from the `text` binding and fill_text_image
# styles TextColumn from `text` as well (the manifest's text_full/text_column
# keys are web-side synonyms bound to the same roles). The emphasis column
# mirrors the renderer's emphasis=True fill sites: only body-flow shapes ever
# carry **keyword** runs; ceremonial single-line roles keep markers literal.
_SHAPE_ROLES = {
    "AgendaLabel": ("label", False),
    "AgendaList": ("list", True),
    "TextArea": ("text", True),
    "TextFullArea": ("text", True),
    "TextColumn": ("text", True),
    "SubheadArea": ("subhead", False),
    "CommentArea": ("comment", True),
}


def _subdict(parent, key: str) -> dict:
    """The child dict at parent[key]; {} for anything missing or malformed —
    manifests are hand-editable (variant manifests are a supported flow), and
    the gate's contract is findings-or-skip, never a traceback."""
    child = parent.get(key) if isinstance(parent, dict) else None
    return child if isinstance(child, dict) else {}


def _sublist(parent, key: str) -> list:
    """The child list at parent[key]; [] unless it is a real list."""
    child = parent.get(key) if isinstance(parent, dict) else None
    return child if isinstance(child, list) else []


def _hex6(color) -> str:
    """"#548235" -> "548235"; already-bare and missing values pass through."""
    return str(color).lstrip("#").upper() if color else ""


def _role_of(sp, bindings: dict, roles: dict, title_text: str | None):
    """(role_name, emphasis_allowed) a shape must render at, None-shaped when
    the archetype binds no usable role for it (stray boxes answer only to the
    sweeps). The title downgrade (44 -> 28 past title_long.over_chars) is
    judged on the rendered title text, like the renderer. Binding values and
    over_chars may be anything a hand-edited manifest left there — non-string
    names and non-numeric thresholds degrade to no expectation."""
    def bound(key):
        name = bindings.get(key)
        return name if isinstance(name, str) else None

    def usable(name):
        return name if isinstance(roles.get(name), dict) else None

    pht = ph_type(sp)
    if pht == "title":
        name = bound("title")
        long_role = roles.get("title_long")
        # over_chars coerces like the renderer and layout.js do (numeric
        # strings parse; anything unparseable means no downgrade)
        try:
            over = float(long_role.get("over_chars", float("inf"))) \
                if isinstance(long_role, dict) else float("inf")
        except (TypeError, ValueError):
            over = float("inf")
        if name == "title" and title_text and vd.text_width(title_text) > over:
            name = "title_long"
        return (usable(name), False)
    if pht == "body":  # cover/closing subtitle; content body phs stay empty
        return (usable(bound("subtitle")), False)
    if sp.name.startswith("Caption"):
        return (usable(bound("caption")), False)
    if sp.name in _SHAPE_ROLES:
        key, emph = _SHAPE_ROLES[sp.name]
        return (usable(bound(key)), emph)
    return None


def _check_style(slide, page, spath: str, manifest: dict) -> list:
    """Tokens at run level: font faces, role styles, math runs, colors.

    Sweeps (font faces, text colors) cover every content run on every slide;
    role expectations need the deck page's archetype and apply only to the
    renderer's named shapes and placeholders. Malformed manifest sections
    degrade to skip, never a traceback.
    """
    tokens = _subdict(_subdict(manifest, "typography"), "tokens")
    if not tokens:
        return []  # a manifest without the token tree: nothing to check against
    roles = _subdict(tokens, "roles")
    emph_cfg = _subdict(tokens, "emphasis")
    emph_color = _hex6(emph_cfg.get("color"))
    emph_bold = emph_cfg.get("weight", "bold") == "bold"  # renderer default
    legal_colors = {_hex6(r.get("color")) for r in roles.values()
                    if isinstance(r, dict) and isinstance(r.get("color"), str)}
    legal_colors |= {emph_color}  # the (bold, emphasis-color) pair is legal
    colors_cfg = _subdict(tokens, "colors")
    accent = _hex6(colors_cfg.get("accent"))
    band = _hex6(colors_cfg.get("band"))
    archetype = page.get("archetype") if isinstance(page, dict) else None
    # closing pages render through fill_cover's cover bindings (render_pptx
    # treats closing as a content variant of cover) — the QA reads what the
    # renderer reads
    binding_arch = "cover" if archetype == "closing" else archetype
    bindings = _subdict(_subdict(tokens, "role_bindings"), binding_arch or "__none__")
    title_text = slide_title_text(slide)
    findings: list[vd.Finding] = []

    for sp in slide.shapes:
        if ph_type(sp) == "sldNum":
            continue  # 页码是母版自带字段（华文中宋），不属于内容样式域
        where = f"{spath}.{sp.name}"
        role_name, emph_ok = _role_of(sp, bindings, roles, title_text) or (None, False)
        role = roles.get(role_name) if role_name else None
        # only a well-formed role (dict with a numeric size_pt) yields a
        # checkable expectation; anything else answers to the sweeps alone
        if not isinstance(role, dict) or not isinstance(role.get("size_pt"), (int, float)):
            role = None
        # face expectation: the role's own face when it overrides (formula;
        # a title-family variant manifest), else the global token family
        face = role.get("face", tokens.get("face")) if role else tokens.get("face")
        latin_face = role.get("face", tokens.get("latin_face")) if role else tokens.get("latin_face")
        for r in sp._element.iter(f"{{{NS_A}}}r"):
            text_el = r.find(f"{{{NS_A}}}t")
            text = (text_el.text or "")[:16] if text_el is not None else ""
            rPr = r.find(f"{{{NS_A}}}rPr")
            latin = rPr.find(f"{{{NS_A}}}latin") if rPr is not None else None
            ea = rPr.find(f"{{{NS_A}}}ea") if rPr is not None else None
            if latin is None or ea is None:
                findings.append(vd.Finding(
                    where, "pptx.font_face",
                    f"run {text!r} lacks an explicit ea/latin typeface; the "
                    f"tokens family is {face!r} (no inheritance)",
                ))
            else:
                for el, want, kind in ((latin, latin_face, "latin"),
                                       (ea, face, "ea")):
                    if el.get("typeface") != want:
                        findings.append(vd.Finding(
                            where, "pptx.font_face",
                            f"run {text!r} {kind} typeface "
                            f"{el.get('typeface')!r} != tokens face {want!r}",
                        ))
            srgb = rPr.find(f"{{{NS_A}}}solidFill/{{{NS_A}}}srgbClr") \
                if rPr is not None else None
            color = _hex6(srgb.get("val")) if srgb is not None else ""
            if color and color not in legal_colors:
                detail = (f"the band color {band} is never a text color"
                          if color == band else
                          f"#{color} is not a tokens role color")
                findings.append(vd.Finding(
                    where, "pptx.text_color", f"run {text!r}: {detail}"))
            if role is not None:
                findings.extend(_check_role_run(
                    rPr, text, role, role_name, emph_ok, emph_color, emph_bold,
                    accent, color, where))
        if sp.name == "FormulaArea":
            formula_role = roles.get("formula")
            if not isinstance(formula_role, dict) or \
                    not isinstance(formula_role.get("size_pt"), (int, float)):
                formula_role = None  # malformed formula role: unverifiable
            findings.extend(_check_math_runs(sp, formula_role, where))
        if sp.name == "TextFullArea":
            # pure-text pages center their paragraphs in the full-height
            # region — a top-anchored sparse page reads as a bottom void
            # (renderer-web mirrors this with flex centering)
            bodyPr = sp._element.find(
                f"{{{NS_P}}}txBody/{{{NS_A}}}bodyPr")
            if bodyPr is None or bodyPr.get("anchor", "t") != "ctr":
                findings.append(vd.Finding(
                    where, "pptx.role_style",
                    "TextFullArea is top-anchored; the text_full design "
                    "centers its paragraphs vertically (anchor='ctr')"))
    return findings


def _check_role_run(rPr, text, role: dict, role_name: str, emph_ok: bool,
                    emph_color: str, emph_bold: bool, accent: str,
                    color: str, where: str) -> list:
    """One run against its role: size always; weight+color as a pair — the
    role pair, or the (bold, emphasis-color) pair inside body-flow shapes.
    The pairing is the doctrine: bold in body flow *is* the emphasis variant,
    and the accent color rides only bold runs. A missing rPr is itself a
    finding — the renderer styles every run explicitly."""
    if rPr is None:
        return [vd.Finding(
            where, "pptx.role_style",
            f"run {text!r} has no rPr; role '{role_name}' expects "
            f"{role['size_pt']}pt {role.get('weight')} {role.get('color')}")]
    findings = []
    want_sz = str(int(role["size_pt"]) * 100)
    if rPr.get("sz") != want_sz:
        findings.append(vd.Finding(
            where, "pptx.role_style",
            f"run {text!r}: font-size rendered {rPr.get('sz')}, role "
            f"'{role_name}' expects {want_sz} ({role['size_pt']}pt)"))
    want_b = "1" if role.get("weight") == "bold" else "0"
    want_color = _hex6(role.get("color"))
    # the emphasis variant mirrors the renderer's own semantics
    # (render_pptx._styled_paragraph): bold only when emphasis.weight says
    # bold, and the emphasis color with a per-role fallback when the token
    # is absent — a variant manifest may legally re-point either
    variants = {(want_b, want_color)}
    if emph_ok:
        variants.add(("1" if emph_bold else want_b, emph_color or want_color))
    got_b = rPr.get("b")
    if (got_b, color) not in variants:
        if not color:
            detail = (f"no explicit solidFill; role '{role_name}' expects "
                      f"{role.get('color')}")
        elif color == accent and want_b == "0":
            detail = (f"accent #{accent} rides only bold runs (emphasis or "
                      f"subhead); role '{role_name}' expects {role.get('color')}")
        elif got_b == "1" and want_b == "0":
            detail = (f"font-weight rendered '1' with color #{color}; bold in "
                      f"body flow is the emphasis variant and pairs with "
                      f"#{emph_color or _hex6(role.get('color'))}")
        else:
            detail = (f"weight/color ({got_b!r}, #{color}) is neither the "
                      f"role '{role_name}' pair ({want_b!r}, {role.get('color')})"
                      f" nor the emphasis variant ('1', #{emph_color or _hex6(role.get('color'))})")
        findings.append(vd.Finding(
            where, "pptx.role_style", f"run {text!r}: {detail}"))
    return findings


def _check_math_runs(sp, formula_role: dict | None, where: str) -> list:
    """Math runs stay on the formula face/size (OMML hard constraint, outside
    the content family token)."""
    if not formula_role:
        return []
    findings = []
    want_face = formula_role.get("face")
    want_sz = str(int(formula_role["size_pt"]) * 100)
    tags = (f"{{{NS_M}}}r", f"{{{NS_M}}}ctrlPr")
    for el in list(sp._element.iter(tags[0])) + list(sp._element.iter(tags[1])):
        rPr = el.find(f"{{{NS_A}}}rPr")
        latin = rPr.find(f"{{{NS_A}}}latin") if rPr is not None else None
        face = latin.get("typeface") if latin is not None else None
        if face != want_face:
            findings.append(vd.Finding(
                where, "pptx.font_face",
                f"math run face {face!r} != formula face {want_face!r}"))
        if rPr is None or rPr.get("sz") != want_sz:
            findings.append(vd.Finding(
                where, "pptx.role_style",
                f"math run size {rPr.get('sz') if rPr is not None else None}, "
                f"formula role expects {want_sz} ({formula_role['size_pt']}pt)"))
    return findings


# brand layer (#20): the master every slide inherits must carry the measured
# brand geometry — bands as shapes at their manifest position, logos as
# pictures whose blip (or svg fallback pair) matches the declared media name.

def _shape_geo(sp):
    try:
        return (Emu(sp.left).inches, Emu(sp.top).inches,
                Emu(sp.width).inches, Emu(sp.height).inches)
    except (TypeError, AttributeError):
        return None


def _geo_close(sp, region: dict, tol: float = 0.05) -> bool:
    geo = _shape_geo(sp)
    return geo is not None and all(
        abs(g - float(region[k])) <= tol
        for g, k in zip(geo, ("x", "y", "w", "h")))


def _pic_media_names(pic, master) -> set[str]:
    """Media file names the picture references: the raster blip plus the SVG
    source when the pic is an svg+png fallback pair."""
    names = set()
    el = pic._element
    for tag, ns in ((f"{{{NS_A}}}blip", NS_R), (f"{{{NS_SVG}}}svgBlip", NS_R)):
        node = el.find(f".//{tag}")
        if node is None:
            continue
        rid = node.get(f"{{{ns}}}embed")
        if rid and rid in master.part.rels:
            names.add(str(master.part.rels[rid].target_ref).rsplit("/", 1)[-1])
    return names


def _region_of(entry) -> dict | None:
    """The x/y/w/h region of a brand_layer entry, None when malformed."""
    if not isinstance(entry, dict):
        return None
    try:
        return {k: float(entry[k]) for k in ("x", "y", "w", "h")}
    except (KeyError, TypeError, ValueError):
        return None


def _check_brand(slide, page, spath: str, manifest: dict) -> list:
    if not isinstance(page, dict):
        return []  # unaligned slide: the archetype (and its brand group) is unknown
    bl = manifest.get("brand_layer")
    bl = bl if isinstance(bl, dict) else {}
    group = _subdict(bl, "cover") if page.get("archetype") in ("cover", "closing") \
        else _subdict(bl, "content")
    if not group:
        return []  # a manifest without brand_layer: the web gate reports it
    master = slide.slide_layout.slide_master
    shapes = list(master.shapes)
    findings = []
    band_key = "mid_band" if "mid_band" in group else "top_band"
    band_region = _region_of(group.get(band_key))
    if band_region is None:
        findings.append(vd.Finding(
            spath, "pptx.brand_layer",
            f"brand_layer group has no usable {band_key} region "
            f"({group.get(band_key)!r})"))
    elif not any(_geo_close(sp, band_region) for sp in shapes):
        findings.append(vd.Finding(
            spath, "pptx.brand_layer",
            f"the master this slide inherits has no {band_key} at the measured "
            f"geometry {band_region}"))
    logos = [(f"logo[{i}]", logo) for i, logo in enumerate(_sublist(group, "logos"))]
    for key in ("corner_logo", "corner_logo_bar"):
        if key in group:
            logos.append((key, group.get(key)))
    for name, logo in logos:
        region = _region_of(logo)
        media = logo.get("media", "") if isinstance(logo, dict) else ""
        if region is None:
            findings.append(vd.Finding(
                spath, "pptx.brand_layer",
                f"brand_layer entry {name} ({logo!r}) is not a usable region"))
            continue
        found = any(
            sp.shape_type == PICTURE_SHAPE_TYPE
            and _geo_close(sp, region)
            and media in _pic_media_names(sp, master)
            for sp in shapes)
        if not found:
            findings.append(vd.Finding(
                spath, "pptx.brand_layer",
                f"the master this slide inherits lacks {name} ({media}) at "
                f"the measured geometry "
                f"({region['x']}, {region['y']}, {region['w']}, {region['h']})"))
    return findings


# ---------------------------------------------------------------------------
# optional COM screenshot spot check (degrades to a stderr note)
# ---------------------------------------------------------------------------


def export_com_screenshots(pptx_path: Path, out_dir: Path) -> tuple[int, str | None]:
    """Export every slide as PNG via PowerPoint COM; (count, error)."""
    try:
        import win32com.client  # noqa: F401
    except ImportError:
        return 0, "pywin32 not installed; COM spot check unavailable"
    try:
        import pythoncom

        pythoncom.CoInitialize()
        import win32com.client

        # DispatchEx: a dedicated instance, so Quit never closes the user's own
        app = win32com.client.DispatchEx("PowerPoint.Application")
    except Exception as exc:
        return 0, f"cannot start PowerPoint COM: {exc}"
    try:
        out_dir.mkdir(parents=True, exist_ok=True)
        pres = app.Presentations.Open(
            str(pptx_path.resolve()), ReadOnly=True, Untitled=False, WithWindow=False)
        try:
            # slide count from the presentation itself: reusing a non-empty
            # DIR must not inflate the reported number with stale PNGs
            n_slides = pres.Slides.Count
            pres.Export(str(out_dir.resolve()), "PNG", 1280, 720)
        finally:
            pres.Close()
        return n_slides, None
    except Exception as exc:
        return 0, f"COM export failed: {exc}"
    finally:
        try:
            app.Quit()
        except Exception:
            pass


# ---------------------------------------------------------------------------
# CLI
# ---------------------------------------------------------------------------


def build_report(findings: list, slide_count: int, fmt: str, extra: dict | None = None) -> str:
    if fmt == "json":
        report = {
            "valid": not findings,
            "slide_count": slide_count,
            "findings": [{"path": f.path, "code": f.code, "message": f.message} for f in findings],
        }
        if extra:
            report.update(extra)
        return json.dumps(report, ensure_ascii=False, indent=2)
    lines = [f"{f.path}: [{f.code}] {f.message}" for f in findings]
    if findings:
        lines.append(f"INVALID: {len(findings)} finding(s)")
    else:
        lines.append(f"OK: {slide_count} slide(s), 0 findings")
    for key, value in (extra or {}).items():
        lines.append(f"{key}: {value}")
    return "\n".join(lines)


def main(argv=None) -> int:
    parser = argparse.ArgumentParser(
        description="QA-check a rendered .pptx against the deck.json it renders.")
    parser.add_argument("pptx", help="path to the rendered .pptx")
    parser.add_argument("--deck", required=True, help="path to the deck.json it renders")
    parser.add_argument("--templates-dir", type=Path,
                        default=Path(__file__).resolve().parent.parent / "templates",
                        help="directory holding <template-id>/ (default: <repo>/templates)")
    parser.add_argument("--manifest", type=Path,
                        help="explicit manifest path (overrides template resolution)")
    parser.add_argument("--com-screenshots", type=Path, metavar="DIR",
                        help="also export slide PNGs via PowerPoint COM into DIR")
    parser.add_argument("--format", choices=("text", "json"), default="text",
                        help="report format (default: text)")
    args = parser.parse_args(argv)
    # gate reports are consumed by programs (and asserted by tests) — keep
    # them UTF-8 regardless of the host locale / pipe encoding
    for stream in (sys.stdout, sys.stderr):
        stream.reconfigure(encoding="utf-8", errors="replace")

    deck_path = Path(args.deck)
    try:
        deck = vd.parse_deck(deck_path.read_text(encoding="utf-8"))
    except (OSError, UnicodeError, ValueError) as exc:
        print(f"error: cannot read deck: {exc}", file=sys.stderr)
        return 2

    manifest_path, error_finding = vd.resolve_manifest_path(
        deck, args.templates_dir, args.manifest)
    if error_finding is not None or manifest_path is None:
        print(build_report([error_finding] if error_finding else [], 0, args.format))
        return 1
    try:
        manifest = json.loads(manifest_path.read_text(encoding="utf-8"))
    except (OSError, ValueError) as exc:
        print(f"error: cannot load manifest {manifest_path}: {exc}", file=sys.stderr)
        return 2

    pptx_path = Path(args.pptx)
    if not pptx_path.is_file():
        print(f"error: pptx not found: {pptx_path}", file=sys.stderr)
        return 2
    try:
        prs = Presentation(pptx_path)
    except Exception as exc:
        finding = vd.Finding("", "pptx.unreadable", f"cannot open with python-pptx: {exc}")
        print(build_report([finding], 0, args.format))
        return 1

    findings = check_pptx(prs, deck, manifest, deck_path.parent)

    extra = {}
    if args.com_screenshots is not None:
        count, error = export_com_screenshots(pptx_path, args.com_screenshots)
        if error is not None:
            print(f"note: {error}; programmatic checks stand alone", file=sys.stderr)
        else:
            extra["com_screenshots"] = f"{count} slide(s) -> {args.com_screenshots}"

    print(build_report(findings, len(prs.slides), args.format, extra))
    return 1 if findings else 0


if __name__ == "__main__":
    sys.exit(main())
