"""Extract template derivatives (brand media, colors, fonts, geometry) from a template pptx.

Usage:
    python scripts/extract_template.py <template.pptx> <extracted_dir>

Writes into <extracted_dir>:
    media/        brand images referenced by slide masters/layouts (original part names)
    media.json    provenance: which master/layout references each part, size, sha256
    colors.json   theme palettes, master->theme map, observed text/fill colors
    fonts.json    measured font table: typeface@size -> runs, chars, slides, role
    geometry.json slide size, master/layout placeholder tables, per-archetype sample slides

The full pptx is never copied; only the derivative files above are written.
"""
import hashlib
import json
import mimetypes
import re
import struct
import sys
import zipfile
from collections import defaultdict
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu

A = "http://schemas.openxmlformats.org/drawingml/2006/main"
P = "http://schemas.openxmlformats.org/presentationml/2006/main"
R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
MATH = "http://schemas.openxmlformats.org/officeDocument/2006/math"


def inches(v):
    return round(Emu(int(v)).inches, 2)


def write_json(path, obj):
    path.write_text(json.dumps(obj, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")


def sha256(data):
    return hashlib.sha256(data).hexdigest()[:16]


def png_dims(data):
    if len(data) > 24 and data[:8] == b"\x89PNG\r\n\x1a\n":
        w, h = struct.unpack(">II", data[16:24])
        return w, h
    return None


def gif_dims(data):
    if len(data) > 10 and data[:6] in (b"GIF87a", b"GIF89a"):
        w, h = struct.unpack("<HH", data[6:10])
        return w, h
    return None


def xfrm_of(sp):
    xfrm = sp.find(f".//{{{A}}}xfrm")
    if xfrm is None:
        return None
    off, ext = xfrm.find(f"{{{A}}}off"), xfrm.find(f"{{{A}}}ext")
    if off is None or ext is None:
        return None
    return {"x": inches(off.get("x")), "y": inches(off.get("y")),
            "w": inches(ext.get("cx")), "h": inches(ext.get("cy"))}


def cNvPr(sp):
    nv = sp.find(f".//{{{P}}}cNvPr")
    return nv.get("name") if nv is not None else ""


def ph_info(sp):
    ph = sp.find(f".//{{{P}}}ph")
    if ph is None:
        return None
    return ph.get("type") or "body", ph.get("idx")


def runs_of(sp):
    out = []
    for para in sp.findall(f".//{{{A}}}p"):
        for run in para.findall(f"{{{A}}}r"):
            rPr = run.find(f"{{{A}}}rPr")
            t = run.find(f"{{{A}}}t")
            text = (t.text or "") if t is not None else ""
            sz = int(rPr.get("sz")) / 100 if rPr is not None and rPr.get("sz") else None
            ea = rPr.find(f"{{{A}}}ea") if rPr is not None else None
            latin = rPr.find(f"{{{A}}}latin") if rPr is not None else None
            color = rPr.find(f".//{{{A}}}srgbClr") if rPr is not None else None
            out.append({
                "ea": ea.get("typeface") if ea is not None else None,
                "latin": latin.get("typeface") if latin is not None else None,
                "sz_pt": sz,
                "chars": len(text),
                "color": color.get("val") if color is not None else None,
                "text": text[:20],
            })
    return out


def text_of(sp):
    return "".join(t.text or "" for t in sp.findall(f".//{{{A}}}t"))


def classify_role(sz):
    if sz is None:
        return "inherit"
    if sz <= 12:
        return "caption"
    if sz >= 40:
        return "display"
    return "body"


def slide_index(prs, slide):
    return list(prs.slides).index(slide) + 1


def dump_shape_tree(part_root):
    """Shapes + pictures with geometry, placeholders, fonts, oMath count."""
    shapes = []
    for sp in part_root.iter(f"{{{P}}}sp"):
        entry = {"kind": "shape", "name": cNvPr(sp), "geo": xfrm_of(sp.find(f"{{{P}}}spPr"))}
        ph = ph_info(sp)
        if ph:
            entry["placeholder"] = {"type": ph[0], "idx": ph[1]}
        runs = runs_of(sp)
        if runs:
            entry["runs"] = runs
        txt = text_of(sp)
        if txt:
            entry["text"] = txt[:40]
        shapes.append(entry)
    for pic in part_root.iter(f"{{{P}}}pic"):
        blip = pic.find(f".//{{{A}}}blip")
        entry = {"kind": "picture", "name": cNvPr(pic), "geo": xfrm_of(pic.find(f"{{{P}}}spPr")),
                 "text": text_of(pic)[:40]}
        if blip is not None:
            entry["rEmbed"] = blip.get(f"{{{R}}}embed")
        shapes.append(entry)
    return shapes


def main(pptx_path, out_dir):
    pptx_path, out_dir = Path(pptx_path), Path(out_dir)
    out_dir.mkdir(parents=True, exist_ok=True)
    prs = Presentation(str(pptx_path))
    zipf = zipfile.ZipFile(str(pptx_path))

    # ---------- media referenced by masters/layouts (brand assets) ----------
    brand_parts = {}  # partname (relative, no leading /) -> {referenced_by: [...]}
    def note_media(part, owner):
        for rel in part.rels.values():
            if "image" in rel.reltype:
                pn = str(rel.target_partname).lstrip("/")
                brand_parts.setdefault(pn, {"referenced_by": []})["referenced_by"].append(owner)

    master_info = []
    theme_of_master = {}
    for mi, master in enumerate(prs.slide_masters, 1):
        owner = f"master{mi}"
        note_media(master.part, owner)
        for layout in master.slide_layouts:
            note_media(layout.part, f"layout:{layout.name}")
        theme_rel = next((r for r in master.part.rels.values() if "theme" in r.reltype), None)
        theme_of_master[owner] = str(theme_rel.target_partname).lstrip("/") if theme_rel else None
        # txStyles (title/body default sizes live here)
        txs = {}
        for style in ("titleStyle", "bodyStyle", "otherStyle"):
            el = master.part._element.find(f".//{{{P}}}{style}")
            if el is not None:
                lvl1 = el.find(f"{{{A}}}lvl1pPr")
                if lvl1 is not None:
                    defrPr = lvl1.find(f"{{{A}}}defRPr")
                    if defrPr is not None:
                        ea = defrPr.find(f"{{{A}}}ea")
                        txs[style] = {"sz_pt": int(defrPr.get("sz")) / 100 if defrPr.get("sz") else None,
                                      "ea": ea.get("typeface") if ea is not None else None}
        master_info.append({
            "master": owner, "theme": theme_of_master[owner],
            "layouts": [{"name": l.name, "placeholders": [
                {"type": p.placeholder_format.type or "body", "idx": p.placeholder_format.idx,
                 "geo": xfrm_of(p._element.find(f"{{{P}}}spPr"))}
                for p in l.placeholders]} for l in master.slide_layouts],
            "txStyles": txs,
        })

    media_dir = out_dir / "media"
    media_dir.mkdir(exist_ok=True)
    media_records = []
    all_media = {n for n in zipf.namelist() if n.startswith("ppt/media/")}
    for pn in sorted(brand_parts):
        data = zipf.read(pn.lstrip("/"))
        fname = pn.rsplit("/", 1)[1]
        (media_dir / fname).write_bytes(data)
        dims = png_dims(data) or gif_dims(data)
        media_records.append({
            "part": pn, "file": f"media/{fname}",
            "content_type": mimetypes.guess_type(fname)[0] or "application/octet-stream",
            "bytes": len(data), "pixel_dims": dims, "sha256_16": sha256(data),
            "referenced_by": sorted(set(brand_parts[pn]["referenced_by"])),
        })
    brand_part_names = {r["part"] for r in media_records}
    content_only = sorted(n.rsplit("/", 1)[1] for n in all_media if n not in brand_part_names)
    write_json(out_dir / "media.json", {
        "brand_media": media_records,
        "excluded_content_media": content_only,
        "note": "excluded_* are figures of the sample presentation (slide content), not template branding",
    })

    # ---------- colors ----------
    def theme_font_scheme(xml, which):
        m = re.search(rf"<a:{which}Font>(.*?)</a:{which}Font>", xml, re.S)
        if not m:
            return None, None
        lat = re.search(r'<a:latin typeface="([^"]*)"', m.group(1))
        ea = re.search(r'<a:ea typeface="([^"]*)"', m.group(1))
        return (lat.group(1) if lat else None, ea.group(1) if ea else None)

    themes = {}
    for n in sorted(zipf.namelist()):
        if re.match(r"ppt/theme/theme\d+\.xml", n):
            xml = zipf.read(n).decode("utf-8")
            name = re.search(r'<a:theme[^>]*name="([^"]*)"', xml)
            pal = re.findall(r'<a:(dk1|lt1|dk2|lt2|accent\d|hlink|folHlink)>.*?(?:val="([0-9A-Fa-f]{6})"|lastClr="([0-9A-Fa-f]{6})")', xml)
            major_l, major_e = theme_font_scheme(xml, "major")
            minor_l, minor_e = theme_font_scheme(xml, "minor")
            themes[n] = {"name": name.group(1) if name else "?",
                         "palette": {c[0]: (c[1] or c[2]).upper() for c in pal},
                         "fontScheme": {"major_latin": major_l, "major_ea": major_e,
                                        "minor_latin": minor_l, "minor_ea": minor_e}}
    text_colors, fill_colors = defaultdict(lambda: {"count": 0, "where": []}), defaultdict(lambda: {"count": 0, "where": []})
    def scan_colors(root, where):
        for rPr in root.iter(f"{{{A}}}rPr"):
            c = rPr.find(f"{{{A}}}solidFill/{{{A}}}srgbClr")
            if c is not None:
                text_colors["#" + c.get("val").upper()]["count"] += 1
                text_colors["#" + c.get("val").upper()]["where"].append(where)
        for spPr in root.iter(f"{{{P}}}spPr"):
            c = spPr.find(f"{{{A}}}solidFill/{{{A}}}srgbClr")
            if c is not None:
                fill_colors["#" + c.get("val").upper()]["count"] += 1
                fill_colors["#" + c.get("val").upper()]["where"].append(where)
    for mi, master in enumerate(prs.slide_masters, 1):
        scan_colors(master.part._element, f"master{mi}")
        for l in master.slide_layouts:
            scan_colors(l.part._element, f"layout:{l.name}")
    for slide in prs.slides:
        scan_colors(slide.part._element, f"slide{slide_index(prs, slide)}")
    for k in list(text_colors):
        text_colors[k]["where"] = sorted(set(text_colors[k]["where"]))
    for k in list(fill_colors):
        fill_colors[k]["where"] = sorted(set(fill_colors[k]["where"]))
    write_json(out_dir / "colors.json", {
        "themes": themes, "master_theme_map": theme_of_master,
        "observed_text_colors": dict(text_colors), "observed_fill_colors": dict(fill_colors),
        "note": "theme palettes are the Office default and theme ea fonts are Arial placeholders; "
                "design identity comes from master images and per-run CJK faces (see fonts.json)",
    })

    # ---------- fonts ----------
    font_rows = defaultdict(lambda: {"runs": 0, "chars": 0, "slides": set(), "roles": set(), "sample": ""})
    def scan_fonts(root, where):
        for rPr in root.iter(f"{{{A}}}rPr"):
            sz = rPr.get("sz")
            if not sz:
                continue
            ea = rPr.find(f"{{{A}}}ea")
            latin = rPr.find(f"{{{A}}}latin")
            face = (ea.get("typeface") if ea is not None else None) or (latin.get("typeface") if latin is not None else None) or "(theme)"
            key = f"{face}@{int(sz)/100:g}pt"
            row = font_rows[key]
            row["runs"] += 1
            row["slides"].add(where)
            run = rPr.getparent()
            txt = ""
            if run is not None:
                tel = run.find(f"{{{A}}}t")
                txt = tel.text or "" if tel is not None else ""
            row["chars"] += len(txt)
            row["roles"].add(classify_role(int(sz) / 100))
            if txt and not row["sample"]:
                row["sample"] = txt[:16]
    for slide in prs.slides:
        scan_fonts(slide.part._element, f"s{slide_index(prs, slide)}")
    fonts_out = {k: {"runs": v["runs"], "chars": v["chars"],
                     "slides": sorted(v["slides"]), "roles": sorted(v["roles"]), "sample": v["sample"]}
                 for k, v in sorted(font_rows.items(), key=lambda kv: -kv[1]["runs"])}
    write_json(out_dir / "fonts.json", {"runs_by_face_size": fonts_out,
                                        "note": "face is ea (CJK) typeface when set, else latin, else theme font"})

    # ---------- geometry ----------
    samples = {}
    for label, idx in [("cover", 1), ("agenda", 2), ("text-formula", 3),
                       ("text-image", 7), ("text-image_alt", 9), ("chart-focus", 18),
                       ("text-full_variant", 11)]:
        slide = list(prs.slides)[idx - 1]
        root = slide.part._element
        samples[label] = {
            "slide": idx, "layout": slide.slide_layout.name,
            "oMath": len(root.findall(f".//{{{MATH}}}oMath")),
            "shapes": dump_shape_tree(root),
        }
    # boxes extending beyond the slide canvas (the template author's own overflow defects)
    sw, sh = inches(prs.slide_width), inches(prs.slide_height)
    overflows = []
    for si, slide in enumerate(prs.slides, 1):
        for shp in dump_shape_tree(slide.part._element):
            g = shp.get("geo")
            if g and (g["x"] < -0.01 or g["y"] < -0.01
                      or g["x"] + g["w"] > sw + 0.01 or g["y"] + g["h"] > sh + 0.01):
                overflows.append({"slide": si, "name": shp["name"], "geo": g})
    slides = list(prs.slides)
    with_notes = [s for s in slides if s.has_notes_slide]
    notes_empty = sum(1 for s in with_notes if not s.notes_slide.notes_text_frame.text.strip())
    write_json(out_dir / "geometry.json", {
        "slide_size": {"emu": [prs.slide_width, prs.slide_height],
                       "inches": [inches(prs.slide_width), inches(prs.slide_height)]},
        "slide_count": len(slides),
        "masters": master_info,
        "archetype_samples": samples,
        "canvas_overflows": overflows,
        "notes_slides_empty": notes_empty,
        "notes_slides_total": len(with_notes),
    })
    print(f"extracted -> {out_dir} ({len(media_records)} brand media, "
          f"{len(fonts_out)} font rows, {len(samples)} samples)")


if __name__ == "__main__":
    if len(sys.argv) != 3:
        sys.exit(__doc__)
    main(sys.argv[1], sys.argv[2])
