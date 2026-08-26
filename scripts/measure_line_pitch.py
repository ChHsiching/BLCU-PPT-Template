"""Measure PowerPoint's rendered line pitch for Noto Sans SC via COM.

Evidence generator for the manifest's capacity model (budget_semantics): the
budget derivation needs the line height PowerPoint actually renders at the
typography.tokens spacing rhythm (spcPct 150%), which depends on the font's
own vertical metrics (Noto Sans SC: ascent 1160 / descent 288 per 1000 upm =
1.448em — far from the 1.2em the historic model assumed). Theory is not
evidence: this script builds one textbox per (size, spcPct) holding a single
10-line paragraph, opens the probe in PowerPoint, reads
TextFrame.TextRange.BoundHeight (points), and reports the per-line pitch and
its em multiple. Run it whenever the font family or the line-spacing token
changes, then update manifest budgets + budget_semantics with the numbers.

Usage:
    python scripts/measure_line_pitch.py [-o probe.pptx] [--face NAME]

Exit codes: 0 measured, 2 environment error (no PowerPoint COM / no pywin32).
"""
from __future__ import annotations

import argparse
import sys
from pathlib import Path

from pptx import Presentation
from pptx.oxml import parse_xml
from pptx.util import Inches, Pt

NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"

# the flowing-text sizes the capacity model budgets (typography.tokens.roles:
# secondary 18 / body 20 / agenda_list 22), measured at the 150% rhythm and at
# 100% so the single-line baseline of the installed face is visible too.
# Keep in sync with the budgeted role sizes when tokens change — the probe
# must re-measure exactly the sizes budgets are derived from.
SIZES_PT = (18, 20, 22)
SPC_PCTS = (100, 150)
LINES = 10
PROBE_TEXT = "组会汇报测试行Padding0123"


def build_probe(path: Path, face: str) -> list[tuple[int, int, int]]:
    """One slide of probe boxes; returns [(size_pt, spc_pct, shape_idx_1based)]."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    probes = []
    top = Inches(0.3)
    for size in SIZES_PT:
        for pct in SPC_PCTS:
            box = slide.shapes.add_textbox(Inches(0.5), top, Inches(12.3), Inches(0.6))
            top += Inches(0.72)
            tf = box.text_frame
            tf.word_wrap = False
            p = tf.paragraphs[0]
            for i in range(LINES):
                run = p.add_run()
                run.text = PROBE_TEXT if i == 0 else "\n" + PROBE_TEXT
                run.font.size = Pt(size)
                run.font.name = face
            pPr = p._p.get_or_add_pPr()
            pPr.append(parse_xml(
                f'<a:lnSpc xmlns:a="{NS_A}"><a:spcPct val="{pct * 1000}"/></a:lnSpc>'))
            probes.append((size, pct, len(slide.shapes)))
    prs.save(path)
    return probes


def measure(path: Path, probes: list[tuple[int, int, int]]) -> dict[tuple[int, int], float]:
    """Read TextRange.BoundHeight per probe box from a live PowerPoint."""
    import win32com.client

    app = win32com.client.Dispatch("PowerPoint.Application")
    pres = app.Presentations.Open(str(path), True, False, False)  # ReadOnly, no window
    try:
        out = {}
        for size, pct, idx in probes:
            shape = pres.Slides(1).Shapes(idx)
            height_pt = shape.TextFrame.TextRange.BoundHeight
            out[(size, pct)] = height_pt / LINES
        return out
    finally:
        pres.Close()
        app.Quit()


def main(argv=None) -> int:
    parser = argparse.ArgumentParser(description=__doc__.splitlines()[0])
    parser.add_argument("-o", "--out", type=Path, default=Path("line-pitch-probe.pptx"))
    parser.add_argument("--face", default="Noto Sans SC")
    args = parser.parse_args(argv)
    for stream in (sys.stdout, sys.stderr):
        stream.reconfigure(errors="replace")

    probes = build_probe(args.out, args.face)
    try:
        pitches = measure(args.out, probes)
    except Exception as exc:
        print(f"error: COM measurement failed: {exc!r}", file=sys.stderr)
        return 2

    print(f"face {args.face!r}, {LINES} lines/box, pitch = BoundHeight/{LINES}")
    for (size, pct), pitch in sorted(pitches.items()):
        print(f"  {size:>2}pt @ {pct:>3}%: {pitch:7.3f} pt/line = {pitch / size:.4f} em")
    return 0


if __name__ == "__main__":
    sys.exit(main())
