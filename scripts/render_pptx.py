"""Render a deck.json to a .pptx by Clone & Fill on the template original.

Usage:
    python scripts/render_pptx.py <deck.json> -o <out.pptx> [--templates-dir DIR] [--manifest PATH]

The template original is opened from an in-memory copy and never written to.
Each deck page is rendered by deep-cloning the archetype's source slide at the
XML level (new slide part, relationships re-registered with rId remapping,
sldIdLst + [Content_Types].xml maintained by python-pptx on save), stripping the
template author's content layer while keeping placeholders and the master brand
layer, then rebuilding content boxes at the manifest regions. Archetypes whose
source slide has overflowing content boxes (agenda, chart-focus) are rebuilt at
the manifest's safe geometry, never at the original boxes' geometry. Formulas
go LaTeX -> MathML -> OMML (MML2OMML.XSL) as native Cambria Math objects in the
same mc:AlternateContent/a14:m serialization the template itself uses; a
formula whose conversion fails falls back to a rendered image inside the
formula region. Images from deck.json blocks are contain-fitted into the
manifest regions/slots (aspect preserved, centered); white-background images
additionally get the manifest image.hairline token (#D9D9D9 0.75pt) so they do
not dissolve into the page. Captions render as 12pt white-on-black-scrim
overlay strips at the owning slot's internal bottom edge. All content
typography is driven by the manifest's typography.tokens role system
(role_bindings resolve regions to roles: Noto Sans SC everywhere, weight +
color hierarchy, **keyword** emphasis runs in green Bold, body spacing rhythm
line 1.5 / before 12pt / 0.1in-0.05in textbox insets), and the rendered fonts
are subset-embedded per the manifest fonts.embed token. Original template
slides are deleted high-index-first after cloning. All 6 archetypes are
supported: cover, agenda, text-formula, text-image, chart-focus, closing.

Exit codes: 0 rendered, 1 deck failed validation, 2 usage / IO / unsupported
archetype / environment error. Decks are validated (scripts/validate_deck.py,
pure function) before anything is cloned.
"""
from __future__ import annotations

import argparse
import copy
import functools
import io
import json
import os
import re
import sys
import zipfile
from dataclasses import dataclass
from pathlib import Path

import lxml.etree as etree

# All XML here (template-derived trees, MathML from LaTeX, the Office XSL,
# this file's own shell templates) must parse with entities unresolved.
_XML_PARSER = etree.XMLParser(resolve_entities=False, no_network=True, load_dtd=False)
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.opc.constants import CONTENT_TYPE as CT
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.oxml import parse_xml
from pptx.oxml.ns import qn
from pptx.parts.slide import SlidePart
from pptx.util import Inches, Pt

sys.path.insert(0, str(Path(__file__).resolve().parent))
import embed_fonts as efont  # noqa: E402
import validate_deck as vd  # noqa: E402

SUPPORTED_ARCHETYPES = (
    "cover", "agenda", "text-formula", "text-image", "chart-focus", "closing",
)

NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
NS_M = "http://schemas.openxmlformats.org/officeDocument/2006/math"
NS_MC = "http://schemas.openxmlformats.org/markup-compatibility/2006"
NS_A14 = "http://schemas.microsoft.com/office/drawing/2010/main"
NS_R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"
NSMAP = {"a": NS_A, "p": NS_P, "r": NS_R, "m": NS_M, "mc": NS_MC, "a14": NS_A14}


def q(tag: str) -> str:
    """'a:r' -> '{drawingml-ns}r' for the namespaces this renderer touches."""
    prefix, local = tag.split(":")
    return f"{{{NSMAP[prefix]}}}{local}"


class RenderError(Exception):
    """Unrecoverable renderer condition (environment, scope, template)."""


# ---------------------------------------------------------------------------
# LaTeX -> MathML -> OMML chain (+ image fallback)
# ---------------------------------------------------------------------------

try:
    import latex2mathml.converter as _l2m

    HAS_LATEX2MATHML = True
except ImportError:  # pragma: no cover - exercised only on missing installs
    _l2m = None
    HAS_LATEX2MATHML = False

_MML2OMML_ENV = "BLCU_MML2OMML_XSL"
_MML2OMML_ROOTS = [
    r"C:\Program Files\Microsoft Office\root",
    r"C:\Program Files\Microsoft Office",
    r"C:\Program Files (x86)\Microsoft Office",
]
_xslt = None


def find_mml2omml_xsl() -> Path | None:
    """Locate Office's MML2OMML.XSL (env override, then Office installs)."""
    env = os.environ.get(_MML2OMML_ENV)
    if env:
        p = Path(env)
        return p if p.is_file() else None
    for root in _MML2OMML_ROOTS:
        hits = sorted(Path(root).rglob("MML2OMML.XSL"))
        if hits:
            return hits[0]
    return None


def _get_xslt(xsl_path: Path | None):
    global _xslt
    if xsl_path is None:
        return None
    if _xslt is None:
        _xslt = etree.XSLT(etree.parse(str(xsl_path), parser=_XML_PARSER))
    return _xslt


def _style_math_runs(omml, size_pt: int) -> None:
    """Give every math run the formula typography from the manifest.

    The XSLT output has no run properties; the template's own math runs carry
    a:rPr with Cambria Math, so we inject sz + italic + Cambria Math
    (latin/ea/cs) into each m:r / m:ctrlPr, after m:rPr when present.
    """
    sz = str(size_pt * 100)
    for el in list(omml.iter()):
        if el.tag not in (q("m:r"), q("m:ctrlPr")):
            continue
        old = el.find(q("a:rPr"))
        if old is not None:
            el.remove(old)
        rPr = etree.Element(q("a:rPr"))
        rPr.set("i", "1")
        rPr.set("sz", sz)
        for tag in ("a:latin", "a:ea", "a:cs"):
            etree.SubElement(rPr, q(tag)).set("typeface", "Cambria Math")
        m_rPr = el.find(q("m:rPr"))
        el.insert(list(el).index(m_rPr) + 1 if m_rPr is not None else 0, rPr)


def convert_latex_to_omml(latex: str, xsl_path: Path | None = None):
    """LaTeX -> MathML -> OMML; returns an m:oMath element or None on failure.

    None means the caller must use the image fallback.
    """
    if not HAS_LATEX2MATHML:
        return None
    try:
        transform = _get_xslt(xsl_path if xsl_path is not None else find_mml2omml_xsl())
        if transform is None:
            return None
        mathml = _l2m.convert(latex)
        result = transform(etree.fromstring(mathml, parser=_XML_PARSER))
        omml = result.getroot()
        if omml.tag != q("m:oMath"):
            return None
        return omml
    except Exception:
        return None


# mathtext covers a large LaTeX subset but not every command latex2mathml
# accepts; normalize the known gaps before handing a formula to the fallback.
_MAThtext_ALIASES = [
    (r"\lVert", r"\Vert"),
    (r"\rVert", r"\Vert"),
    (r"\lvert", r"\vert"),
    (r"\rvert", r"\vert"),
]


def render_latex_png_bytes(latex: str) -> bytes:
    """Fallback: render LaTeX with matplotlib mathtext to a transparent PNG."""
    try:
        import matplotlib

        matplotlib.use("Agg")
        import matplotlib.pyplot as plt

        tex = latex
        for old, new in _MAThtext_ALIASES:
            tex = tex.replace(old, new)
        fig = plt.figure(figsize=(8, 1.2))
        fig.text(0.5, 0.5, f"${tex}$", ha="center", va="center", fontsize=14)
        buf = io.BytesIO()
        fig.savefig(buf, format="png", dpi=300, transparent=True,
                    bbox_inches="tight", pad_inches=0.02)
        plt.close(fig)
        return buf.getvalue()
    except Exception as exc:  # mathtext cannot parse, or matplotlib missing
        raise RenderError(f"formula fallback rendering failed for: {latex!r} ({exc})") from exc


# ---------------------------------------------------------------------------
# typography tokens + XML building blocks (element level; mirror the
# template's own serialization)
# ---------------------------------------------------------------------------

def _srgb_val(color: str) -> str:
    """"#548235" -> "548235" for a:srgbClr/@val."""
    return color.lstrip("#").upper()


def _tokens(manifest: dict) -> dict:
    """The S2 token tree (typography.tokens) — the style single source of truth."""
    return manifest["typography"]["tokens"]


def _role_style(tokens: dict, role_name: str) -> dict:
    """Resolve a named role to concrete run attrs; face/latin inherit the
    token family defaults (only the formula role carries its own face, and
    math runs are styled by _style_math_runs, never through here)."""
    role = tokens["roles"][role_name]
    return {
        "size_pt": role["size_pt"],
        "face": role.get("face", tokens["face"]),
        "latin": role.get("face", tokens["latin_face"]),
        "bold": role["weight"] == "bold",
        "color": role["color"],
    }


def _region_role_name(tokens: dict, archetype: str, region: str) -> str:
    """The one path through role_bindings: region -> role name."""
    return tokens["role_bindings"][archetype][region]


def _region_role(tokens: dict, archetype: str, region: str) -> dict:
    """Region -> role style via typography.tokens.role_bindings."""
    return _role_style(tokens, _region_role_name(tokens, archetype, region))


def _title_style(text: str, tokens: dict, archetype: str) -> dict:
    """Title role for this page; past title_long.over_chars (CJK width) the
    long-title downgrade applies (44 -> 28 Bold, weight/color unchanged)."""
    name = _region_role_name(tokens, archetype, "title")
    long_role = tokens["roles"].get("title_long")
    if name == "title" and long_role and \
            vd.text_width(text) > long_role.get("over_chars", float("inf")):
        name = "title_long"
    return _role_style(tokens, name)


@functools.lru_cache(maxsize=None)
def _emphasis_pattern(marker: str) -> re.Pattern:
    """The **emphasis** marker regex (paired, non-greedy, DOTALL); markers
    around empty or unpaired content never match and stay literal. Cached —
    the marker is fixed per manifest, so it compiles once per process, not
    once per paragraph."""
    return re.compile(re.escape(marker) + r"(.+?)" + re.escape(marker), re.S)


def _split_emphasis(text: str, marker: str) -> list[tuple[str, bool]]:
    """Split into (segment, emphasized) pairs; unpaired markers stay literal."""
    if marker and marker in text:
        # one capture group -> [literal, captured, literal, ...]; the captured
        # segments (odd indices) are the emphasized ones
        parts = _emphasis_pattern(marker).split(text)
        return [(seg, i % 2 == 1) for i, seg in enumerate(parts) if seg]
    return [(text, False)]


def _styled_paragraph(text: str, style: dict, tokens: dict, *,
                      algn: str | None = None, rhythm: bool = False,
                      emphasis: bool = False) -> etree._Element:
    """One paragraph in a resolved role style.

    rhythm=True applies the body spacing tokens (line 1.5 / before 12pt from
    typography.tokens.spacing) to flowing body text; ceremonial single-line
    roles stay single-spaced with no space-before. emphasis=True (body-flow
    roles only, per the 正文内 doctrine) turns **keyword** segments into
    emphasis runs (tokens.emphasis: bold + accent green); every other role
    renders markers literally.
    """
    spacing = tokens["spacing"]
    ln_pct = int(float(spacing["line_spacing"]) * 100) if rhythm else 100
    spc_bef = float(spacing["space_before_pt"]) if rhythm else 0
    emph_cfg = tokens.get("emphasis", {})
    marker = emph_cfg.get("marker", "")
    emph_bold = emph_cfg.get("weight", "bold") == "bold"
    emph_color = emph_cfg.get("color", style["color"])

    algn_attr = f' algn="{algn}"' if algn else ""
    pPr = f'<a:pPr indent="0"{algn_attr}>'
    pPr += f'<a:lnSpc><a:spcPct val="{ln_pct * 1000}"/></a:lnSpc>'
    if spc_bef:
        pPr += f'<a:spcBef><a:spcPts val="{int(spc_bef * 100)}"/></a:spcBef>'
    pPr += '<a:buNone/></a:pPr>'

    def run_xml(seg: str, bold: bool, color: str) -> str:
        b = "1" if bold else "0"
        rPr = (f'<a:rPr lang="zh-CN" sz="{style["size_pt"] * 100}" b="{b}"'
               f' strike="noStrike" spc="-1">'
               f'<a:solidFill><a:srgbClr val="{_srgb_val(color)}"/></a:solidFill>'
               f'<a:latin typeface="{style["latin"]}"/>'
               f'<a:ea typeface="{style["face"]}"/></a:rPr>')
        escaped = (seg.replace("&", "&amp;").replace("<", "&lt;")
                   .replace(">", "&gt;"))
        return f"<a:r>{rPr}<a:t>{escaped}</a:t></a:r>"

    segments = _split_emphasis(text, marker) if emphasis else [(text, False)]
    runs = "".join(
        run_xml(seg, style["bold"] or (is_emph and emph_bold),
                emph_color if is_emph else style["color"])
        for seg, is_emph in segments)
    end = (f'<a:endParaRPr lang="en-US" sz="{style["size_pt"] * 100}"'
           f' b="{"1" if style["bold"] else "0"}" strike="noStrike" spc="-1">'
           f'<a:solidFill><a:srgbClr val="{_srgb_val(style["color"])}"/></a:solidFill>'
           f"</a:endParaRPr>")
    xml = f'<a:p xmlns:a="{NS_A}">{pPr}{runs}{end}</a:p>'
    # parse_xml (python-pptx's parser) attaches the CT_TextParagraph class so
    # the paragraph behaves inside python-pptx object trees; input is this
    # module's own template with escaped text, never external content.
    return parse_xml(xml)


_MATH_PARA_SHELL = (
    f'<a:p xmlns:a="{NS_A}" xmlns:mc="{NS_MC}" xmlns:a14="{NS_A14}" xmlns:m="{NS_M}">'
    # 公式段统一节奏：居中对齐页面中轴、130% 行距、段前 16pt——
    # 消除宽度不等的公式左缘参差与垂直间距忽紧忽松
    '<a:pPr algn="ctr"><a:lnSpc><a:spcPct val="130000"/></a:lnSpc>'
    '<a:spcBef><a:spcPts val="1600"/></a:spcBef><a:buNone/></a:pPr>'
    "<mc:AlternateContent>"
    '<mc:Choice Requires="a14"><a14:m><m:oMathPara>'
    '<m:oMathParaPr><m:jc m:val="centerGroup"/></m:oMathParaPr>'
    "</m:oMathPara></a14:m></mc:Choice>"
    "<mc:Fallback/>"
    "</mc:AlternateContent>"
    '<a:endParaRPr lang="en-US"/>'
    "</a:p>"
)


def _math_paragraph(omml) -> etree._Element:
    """Wrap an m:oMath the way the template serializes slide math."""
    p = parse_xml(_MATH_PARA_SHELL)
    holder = p.find(f".//{q('m:oMathPara')}")
    holder.append(omml)
    return p


def _textbox_shape(name: str, region: dict, shape_id: int, tokens: dict,
                   anchor: str = "t", fill_xml: str | None = None) -> etree._Element:
    x, y = int(Inches(region["x"])), int(Inches(region["y"]))
    cx, cy = int(Inches(region["w"])), int(Inches(region["h"]))
    spacing = tokens["spacing"]
    h_in = int(Inches(float(spacing["textbox_inset_in"])))
    v_in = int(Inches(float(spacing["textbox_inset_v_in"])))
    fill = fill_xml if fill_xml is not None else "<a:noFill/>"
    xml = (
        f'<p:sp xmlns:a="{NS_A}" xmlns:p="{NS_P}">'
        "<p:nvSpPr>"
        f'<p:cNvPr id="{shape_id}" name="{name}"/><p:cNvSpPr txBox="1"/><p:nvPr/>'
        "</p:nvSpPr>"
        '<p:spPr bwMode="auto">'
        f'<a:xfrm><a:off x="{x}" y="{y}"/><a:ext cx="{cx}" cy="{cy}"/></a:xfrm>'
        '<a:prstGeom prst="rect"><a:avLst/></a:prstGeom>'
        f"{fill}<a:ln w=\"0\"><a:noFill/></a:ln>"
        "</p:spPr>"
        "<p:txBody>"
        f'<a:bodyPr lIns="{h_in}" tIns="{v_in}" rIns="{h_in}" bIns="{v_in}" '
        f'anchor="{anchor}" wrap="square">'
        "<a:noAutofit/></a:bodyPr>"
        "<a:lstStyle/>"
        "</p:txBody>"
        "</p:sp>"
    )
    # parse_xml: CT_Shape class attaches, so slide.shapes iteration works
    return parse_xml(xml)


def _sp_tree(slide_el):
    return slide_el.find(q("p:cSld")).find(q("p:spTree"))


def _next_shape_id(slide_el) -> int:
    ids = [int(c.get("id")) for c in slide_el.iter(q("p:cNvPr")) if c.get("id")]
    return max(ids, default=1) + 1


def ph_type(sp_el) -> str | None:
    """Placeholder type of a p:sp, None when not a placeholder.

    A bare <p:ph> (no type attribute) is a body placeholder.
    """
    ph = sp_el.find(f".//{q('p:ph')}")
    if ph is None:
        return None
    return ph.get("type") or "body"


def find_placeholder(slide_el, want: str):
    for sp in _sp_tree(slide_el).findall(q("p:sp")):
        if ph_type(sp) == want:
            return sp
    return None


def _set_position(sp_el, region: dict) -> None:
    xfrm = sp_el.find(f"{q('p:spPr')}/{q('a:xfrm')}")
    off, ext = xfrm.find(q("a:off")), xfrm.find(q("a:ext"))
    off.set("x", str(int(Inches(region["x"]))))
    off.set("y", str(int(Inches(region["y"]))))
    ext.set("cx", str(int(Inches(region["w"]))))
    ext.set("cy", str(int(Inches(region["h"]))))


def _fill_textbox(sp_el, paragraphs) -> None:
    txBody = sp_el.find(q("p:txBody"))
    for p_el in txBody.findall(q("a:p")):
        txBody.remove(p_el)
    for p_el in paragraphs:
        txBody.append(p_el)


def _fill_placeholder(sp_el, text: str, style: dict, tokens: dict,
                      algn: str | None = None) -> None:
    _fill_textbox(sp_el, [_styled_paragraph(text, style, tokens, algn=algn)])


# ---------------------------------------------------------------------------
# images (deck-referenced material files), hairline, and captions
# ---------------------------------------------------------------------------


def _is_white_backgrounded(image_path: Path) -> bool:
    """True when the image's border is uniformly near-white (screenshot look).
    Degrades to False without PIL: no hairline rather than a wrong one."""
    try:
        from PIL import Image
    except ImportError:
        return False
    try:
        with Image.open(image_path) as im:
            rgb = im.convert("RGB")
            w, h = rgb.size
            points = [(0, 0), (w - 1, 0), (0, h - 1), (w - 1, h - 1),
                      (w // 2, 0), (w // 2, h - 1), (0, h // 2), (w - 1, h // 2)]
            return all(min(rgb.getpixel(p)) >= 245 for p in points)
    except Exception:  # unreadable image: the fit already failed if corrupt
        return False


def _apply_hairline(pic, image_path: Path, tokens: dict) -> None:
    """White-background screenshots get the image.hairline token so they do
    not dissolve into the white page; photos/dark images get nothing."""
    hair = tokens.get("image", {}).get("hairline")
    if not hair or not _is_white_backgrounded(image_path):
        return
    pic.line.color.rgb = RGBColor.from_string(_srgb_val(hair["color"]))
    pic.line.width = Pt(float(hair["width_pt"]))


def _add_fitted_picture(slide, image_path: Path, slot: dict, tokens: dict):
    """Contain-fit an image into a manifest slot: aspect kept, centered."""
    image_path = Path(image_path)
    if not image_path.is_file():
        raise RenderError(f"image file not found: {image_path}")
    left, top = Inches(slot["x"]), Inches(slot["y"])
    slot_w, slot_h = Inches(slot["w"]), Inches(slot["h"])
    with open(image_path, "rb") as stream:
        pic = slide.shapes.add_picture(stream, left, top, width=slot_w)
    if pic.height > slot_h:  # width-bound fit overflows: rebind to height
        pic.width = int(pic.width * (slot_h / pic.height))
        pic.height = slot_h
    pic.left = int(left + (slot_w - pic.width) / 2)
    pic.top = int(top + (slot_h - pic.height) / 2)
    _apply_hairline(pic, image_path, tokens)
    return pic


def _scrim_fill_xml(scrim: dict) -> str:
    return (f'<a:solidFill><a:srgbClr val="{_srgb_val(scrim["color"])}">'
            f'<a:alpha val="{int(scrim["alpha_pct"] * 1000)}"/></a:srgbClr>'
            f"</a:solidFill>")


def _caption_strip_height_in(tokens: dict, caption_role: str) -> float:
    """One caption line at single spacing (PowerPoint renders spcPct 100% as
    1.2em — measured, scripts/measure_line_pitch.py) plus both v-insets."""
    size_pt = tokens["roles"][caption_role]["size_pt"]
    return size_pt * 1.2 / 72 + 2 * float(tokens["spacing"]["textbox_inset_v_in"])


def _add_caption(slide_el, pic, text: str, tokens: dict, archetype: str,
                 index: int) -> None:
    """White-on-black-scrim strip over the fitted picture's bottom edge.

    The caption role (text style + scrim) resolves through the owning
    archetype's role_binding. The strip hugs the fitted image, not the slot:
    a contain-fitted picture may be narrower/shorter than its slot, and a
    scrim floating past the image edges (or below a width-bound image)
    reads as a defect.
    """
    role_name = _region_role_name(tokens, archetype, "caption")
    style = _role_style(tokens, role_name)
    strip_h = _caption_strip_height_in(tokens, role_name)
    strip = {
        "x": pic.left.inches,
        "y": pic.top.inches + pic.height.inches - strip_h,
        "w": pic.width.inches,
        "h": strip_h,
    }
    scrim = tokens["roles"][role_name]["scrim"]
    box = _textbox_shape(f"Caption{index}", strip, _next_shape_id(slide_el),
                         tokens, anchor="ctr", fill_xml=_scrim_fill_xml(scrim))
    _fill_textbox(box, [_styled_paragraph(text, style, tokens, algn="ctr")])
    _sp_tree(slide_el).append(box)


def _strip_content_shapes(slide_el) -> None:
    """Remove the template author's content layer; placeholders stay."""
    tree = _sp_tree(slide_el)
    for child in list(tree):
        if child.tag == q("p:nvGrpSpPr") or child.tag == q("p:grpSpPr"):
            continue
        if child.tag == q("p:sp") and ph_type(child) is not None:
            continue
        tree.remove(child)


# ---------------------------------------------------------------------------
# slide cloning (part + rels + rId remap) and original deletion
# ---------------------------------------------------------------------------

def clone_slide(prs, source_slide):
    """Deep-clone a slide: new part, rels re-registered, rIds remapped, appended."""
    src_part = source_slide.part
    package = prs.part.package
    partname = package.next_partname("/ppt/slides/slide%d.xml")
    element = copy.deepcopy(src_part._element)
    new_part = SlidePart(partname, CT.PML_SLIDE, package, element)

    rid_map = {}
    for old_rid in sorted(src_part.rels.keys()):
        rel = src_part.rels[old_rid]
        if rel.reltype == RT.NOTES_SLIDE:
            continue  # speaker notes never travel with cloned slides
        if rel.is_external:
            new_rid = new_part.rels.get_or_add_ext_rel(rel.reltype, rel.target_ref)
        else:
            new_rid = new_part.rels.get_or_add(rel.reltype, rel.target_part)
        rid_map[old_rid] = new_rid

    if rid_map:
        for el in element.iter():
            for attr, value in list(el.attrib.items()):
                if attr.startswith(f"{{{NS_R}}}") and value in rid_map:
                    el.set(attr, rid_map[value])

    r_id = prs.part.relate_to(new_part, RT.SLIDE)
    prs.slides._sldIdLst.add_sldId(r_id)
    return new_part.slide


def _delete_leading_slides(prs, count: int) -> None:
    """Drop the original template slides, high index first."""
    entries = list(prs.slides._sldIdLst)
    for sldId in reversed(entries[:count]):
        prs.part.drop_rel(sldId.get(qn("r:id")))
        prs.slides._sldIdLst.remove(sldId)


# ---------------------------------------------------------------------------
# archetype fillers
# ---------------------------------------------------------------------------

def _page_texts(page: dict) -> list[str]:
    return [b["text"] for b in page["blocks"] if b["type"] == "text"]


def _page_title(page: dict) -> str:
    return next(b["text"] for b in page["blocks"] if b["type"] == "title")


def _fill_title(slide_el, page: dict, arch: dict, tokens: dict, archetype: str,
                algn: str | None = None) -> None:
    """Fill the title placeholder of a cloned slide at the manifest region."""
    title_text = _page_title(page)
    title = find_placeholder(slide_el, "title")
    if title is None:
        raise RenderError(f"{archetype} clone lost its title placeholder")
    _set_position(title, arch["regions"]["title"])
    _fill_placeholder(title, title_text,
                      _title_style(title_text, tokens, archetype),
                      tokens, algn=algn)


def fill_cover(slide, page: dict, arch: dict, manifest: dict) -> None:
    el = slide.part._element
    tokens = _tokens(manifest)
    text = " ".join(_page_texts(page))

    _strip_content_shapes(el)
    _fill_title(el, page, arch, tokens, "cover", algn="ctr")
    subtitle = find_placeholder(el, "body")
    if subtitle is None:
        raise RenderError("cover clone lost its subtitle placeholder")
    _set_position(subtitle, arch["regions"]["subtitle"])
    _fill_placeholder(subtitle, text, _region_role(tokens, "cover", "subtitle"),
                      tokens, algn="ctr")


def fill_text_formula(slide, page: dict, arch: dict, manifest: dict,
                      xsl_path: Path | None, stats) -> None:
    el = slide.part._element
    tokens = _tokens(manifest)
    formulas = [b["latex"] for b in page["blocks"] if b["type"] == "formula"]
    texts = _page_texts(page)
    formula_size = tokens["roles"]["formula"]["size_pt"]
    body_style = _region_role(tokens, "text-formula", "text")

    _strip_content_shapes(el)
    _fill_title(el, page, arch, tokens, "text-formula")

    # formula area: native OMML paragraphs, image fallbacks stacked below
    native, fallbacks = [], []
    for latex in formulas:
        omml = convert_latex_to_omml(latex, xsl_path)
        if omml is None:
            fallbacks.append(latex)
        else:
            _style_math_runs(omml, formula_size)
            native.append(_math_paragraph(omml))
    stats.formulas_native += len(native)
    stats.formulas_fallback += len(fallbacks)

    if native:
        box = _textbox_shape("FormulaArea", arch["regions"]["formula"],
                             _next_shape_id(el), tokens)
        _fill_textbox(box, native)
        _sp_tree(el).append(box)
    # Fallback images live in the formula region's usable band: clamped above
    # the text region (the manifest lets the formula and text regions overlap;
    # native flow never enters that band, images must not either) and below
    # the headroom reserved for native paragraphs flowing from the top.
    for i, latex in enumerate(fallbacks):
        png = io.BytesIO(render_latex_png_bytes(latex))
        region = arch["regions"]["formula"]
        band_top, band_bottom = Inches(region["y"]), Inches(region["y"]) + Inches(region["h"])
        if texts:
            band_bottom = min(band_bottom, Inches(arch["regions"]["text"]["y"]))
        if native:
            band_top += min(Inches(0.55) * len(native), (band_bottom - band_top) // 2)
        slot = (band_bottom - band_top) / len(fallbacks)
        height = int(slot * 0.9)
        top = int(band_bottom - slot * (i + 1) + (slot - height) / 2)
        pic = slide.shapes.add_picture(png, Inches(region["x"]), top, height=height)
        pic.left = int(Inches(region["x"]) + (Inches(region["w"]) - pic.width) / 2)

    if texts:
        # role_bindings map both text and text_full to the body role; the
        # region (and box name) depends on whether formulas share the page
        region = arch["regions"]["text"] if formulas else arch["regions"]["text_full"]
        paras = [_styled_paragraph(t, body_style, tokens, rhythm=True,
                                   emphasis=True)
                 for t in texts]
        box = _textbox_shape("TextArea" if formulas else "TextFullArea",
                             region, _next_shape_id(el), tokens)
        _fill_textbox(box, paras)
        _sp_tree(el).append(box)


def fill_agenda(slide, page: dict, arch: dict, manifest: dict) -> None:
    """Rebuild label + list at the manifest's safe geometry.

    The agenda source (slide 2) has no title placeholder and both of its
    original content boxes overflow the canvas; the page's title block fills
    the centered label box instead.
    """
    el = slide.part._element
    tokens = _tokens(manifest)
    label_text = _page_title(page)
    items = [item for b in page["blocks"] if b["type"] == "list" for item in b["items"]]
    label_style = _region_role(tokens, "agenda", "label")
    list_style = _region_role(tokens, "agenda", "list")

    _strip_content_shapes(el)
    label = _textbox_shape("AgendaLabel", arch["regions"]["label"],
                           _next_shape_id(el), tokens, anchor="ctr")
    _fill_textbox(label, [_styled_paragraph(label_text, label_style, tokens,
                                            algn="ctr")])
    _sp_tree(el).append(label)
    if items:  # a list block is optional; never emit an empty txBody
        box = _textbox_shape("AgendaList", arch["regions"]["list"],
                             _next_shape_id(el), tokens)
        _fill_textbox(box, [_styled_paragraph(item, list_style, tokens,
                                              rhythm=True, emphasis=True)
                            for item in items])
        _sp_tree(el).append(box)


def fill_text_image(slide, page: dict, arch: dict, manifest: dict,
                    image_root: Path | None) -> None:
    el = slide.part._element
    tokens = _tokens(manifest)
    subhead = next((b["text"] for b in page["blocks"] if b["type"] == "subhead"), None)
    subhead_style = _region_role(tokens, "text-image", "subhead")
    text_style = _region_role(tokens, "text-image", "text")

    _strip_content_shapes(el)
    _fill_title(el, page, arch, tokens, "text-image")

    if subhead is not None:
        box = _textbox_shape("SubheadArea", arch["regions"]["subhead"],
                             _next_shape_id(el), tokens)
        _fill_textbox(box, [_styled_paragraph(subhead, subhead_style, tokens)])
        _sp_tree(el).append(box)

    # text blocks and list items share the text region; the template author
    # writes bullet lines as "- " paragraphs in that box (slide 7)
    paragraphs = []
    for b in page["blocks"]:
        if b["type"] == "text":
            paragraphs.append(_styled_paragraph(b["text"], text_style, tokens,
                                                rhythm=True, emphasis=True))
        elif b["type"] == "list":
            for item in b["items"]:
                paragraphs.append(_styled_paragraph(f"- {item}", text_style,
                                                    tokens, rhythm=True,
                                                    emphasis=True))

    images = [b for b in page["blocks"] if b["type"] == "image"]
    # 单图变体：左文右图整版（设计布局，非模板原框）——图占右栏大区，
    # 文字入左栏纵列；多图仍走模板原 slots 顺序摆位
    single = len(images) == 1 and "image_primary" in arch["regions"]
    if single:
        if paragraphs:
            box = _textbox_shape("TextColumn", arch["regions"]["text_column"],
                                 _next_shape_id(el), tokens)
            _fill_textbox(box, paragraphs)
            _sp_tree(el).append(box)
        slot = arch["regions"]["image_primary"]
        pic = _add_fitted_picture(
            slide, vd.resolve_image_path(images[0]["path"], image_root),
            slot, tokens)
        if images[0].get("caption"):
            _add_caption(el, pic, images[0]["caption"], tokens, "text-image", 1)
        return

    if paragraphs:
        box = _textbox_shape("TextArea", arch["regions"]["text"], _next_shape_id(el),
                             tokens)
        _fill_textbox(box, paragraphs)
        _sp_tree(el).append(box)

    slots = arch["regions"]["image_slots"]
    if len(images) > len(slots):
        raise RenderError(
            f"text-image page carries {len(images)} images but the manifest has "
            f"only {len(slots)} image slots")
    for i, block in enumerate(images):
        pic = _add_fitted_picture(
            slide, vd.resolve_image_path(block["path"], image_root),
            slots[i], tokens)
        if block.get("caption"):
            _add_caption(el, pic, block["caption"], tokens, "text-image", i + 1)


def fill_chart_focus(slide, page: dict, arch: dict, manifest: dict,
                     image_root: Path | None) -> None:
    el = slide.part._element
    tokens = _tokens(manifest)
    texts = _page_texts(page)
    images = [b for b in page["blocks"] if b["type"] == "image"]
    comment_style = _region_role(tokens, "chart-focus", "comment")

    _strip_content_shapes(el)
    _fill_title(el, page, arch, tokens, "chart-focus")

    if len(images) != 1:
        raise RenderError(f"chart-focus expects exactly 1 image block, got {len(images)}")
    _add_fitted_picture(slide, vd.resolve_image_path(images[0]["path"], image_root),
                        arch["regions"]["chart"], tokens)

    if texts:
        box = _textbox_shape("CommentArea", arch["regions"]["comment"],
                             _next_shape_id(el), tokens)
        _fill_textbox(box, [_styled_paragraph(t, comment_style, tokens,
                                              rhythm=True, emphasis=True)
                            for t in texts])
        _sp_tree(el).append(box)


# ---------------------------------------------------------------------------
# pipeline
# ---------------------------------------------------------------------------

@dataclass
class RenderStats:
    pages: int = 0
    formulas_native: int = 0
    formulas_fallback: int = 0
    # fonts.embed payload from embed_fonts (None when not requested/degraded)
    fonts_embedded: dict | None = None
    # degrade reason when embedding was requested but skipped
    font_warning: str | None = None


@dataclass
class RenderResult:
    presentation: object
    stats: RenderStats


def render_deck(deck: dict, manifest: dict, template_pptx: Path,
                xsl_path: Path | None = None,
                image_root: Path | None = None) -> RenderResult:
    """Render a validated deck dict against a template manifest.

    image_root is the directory deck image paths resolve against (the deck
    file's directory at the CLI layer), mirroring validate_deck.
    """
    for i, page in enumerate(deck.get("pages", [])):
        if page.get("archetype") not in SUPPORTED_ARCHETYPES:
            raise RenderError(
                f"pages[{i}]: archetype '{page.get('archetype')}' is not implemented "
                f"in renderer-pptx (supported: {', '.join(SUPPORTED_ARCHETYPES)})"
            )

    template_pptx = Path(template_pptx)
    if not template_pptx.is_file():
        raise RenderError(f"template original not found: {template_pptx}")

    prs = Presentation(io.BytesIO(template_pptx.read_bytes()))
    original_count = len(prs.slides._sldIdLst)
    template_slides = list(prs.slides)
    sources = {}
    for name in SUPPORTED_ARCHETYPES:
        n = manifest["archetypes"][name]["clone"]["slide"]
        sources[name] = template_slides[n - 1]

    if xsl_path is None:
        xsl_path = find_mml2omml_xsl()

    stats = RenderStats()
    for page in deck["pages"]:
        slide = clone_slide(prs, sources[page["archetype"]])
        arch = manifest["archetypes"][page["archetype"]]
        name = page["archetype"]
        if name in ("cover", "closing"):  # closing is a content variant of cover
            fill_cover(slide, page, arch, manifest)
        elif name == "agenda":
            fill_agenda(slide, page, arch, manifest)
        elif name == "text-formula":
            fill_text_formula(slide, page, arch, manifest, xsl_path, stats)
        elif name == "text-image":
            fill_text_image(slide, page, arch, manifest, image_root)
        else:  # chart-focus; scope loop above guarantees the closed set
            fill_chart_focus(slide, page, arch, manifest, image_root)
        stats.pages += 1

    _delete_leading_slides(prs, original_count)

    # font embedding (manifest fonts.embed; degrade with warning, never fail).
    # Characters are collected raw: emphasis markers stay in because literal
    # roles (caption, title) render them verbatim — a superset subset costs a
    # byte, an under-subset breaks machines without the font.
    embed_cfg = manifest.get("fonts", {}).get("embed")
    if embed_cfg:
        base = Path(template_pptx).parent
        try:
            stats.fonts_embedded = efont.embed_fonts(
                prs, face=embed_cfg["face"],
                weights={name: base / embed_cfg[name]
                         for name in efont.WEIGHT_ELEMENTS if name in embed_cfg},
                characters=efont.collect_deck_characters(deck),
                warn=lambda msg: setattr(stats, "font_warning", msg))
        except Exception as exc:  # embedding is an enhancement, never fatal
            stats.font_warning = f"font embedding failed: {exc!r}"
    return RenderResult(prs, stats)


# ---------------------------------------------------------------------------
# CLI
# ---------------------------------------------------------------------------

def main(argv=None) -> int:
    parser = argparse.ArgumentParser(
        description="Render a deck.json to .pptx via Clone & Fill on the template original.")
    parser.add_argument("deck", help="path to deck.json")
    parser.add_argument("-o", "--out", required=True, help="output .pptx path")
    parser.add_argument("--templates-dir", type=Path,
                        default=Path(__file__).resolve().parent.parent / "templates",
                        help="directory holding <template-id>/ (default: <repo>/templates)")
    parser.add_argument("--manifest", type=Path,
                        help="explicit manifest path (overrides template resolution)")
    args = parser.parse_args(argv)
    for stream in (sys.stdout, sys.stderr):
        stream.reconfigure(errors="replace")

    deck_path = Path(args.deck)
    try:
        raw = deck_path.read_text(encoding="utf-8")
        deck = vd.parse_deck(raw)
    except (OSError, UnicodeError) as exc:
        print(f"error: cannot read deck: {exc}", file=sys.stderr)
        return 2
    except ValueError as exc:
        print(f"error: deck is not valid JSON: {exc}", file=sys.stderr)
        return 2
    deck_dir = deck_path.parent

    manifest_path, error_finding = vd.resolve_manifest_path(
        deck, args.templates_dir, args.manifest)
    if error_finding is not None or manifest_path is None:
        print(vd.build_report([error_finding] if error_finding else [], 0, "text"))
        return 1
    try:
        manifest = json.loads(manifest_path.read_text(encoding="utf-8"))
    except (OSError, ValueError) as exc:
        print(f"error: cannot load manifest {manifest_path}: {exc}", file=sys.stderr)
        return 2

    findings = vd.validate_deck(deck, manifest, deck_dir)
    if findings:
        print(vd.build_report(findings, len(deck["pages"]), "text"))
        return 1

    template_pptx = manifest_path.parent / manifest.get("source_pptx", "")
    if not manifest.get("source_pptx"):
        print(f"error: manifest {manifest_path} lacks 'source_pptx'", file=sys.stderr)
        return 2
    if os.path.normcase(str(Path(args.out).resolve())) == \
            os.path.normcase(str(template_pptx.resolve())):
        print("error: refusing to overwrite the template original: "
              f"{template_pptx}", file=sys.stderr)
        return 2
    try:
        result = render_deck(deck, manifest, template_pptx, image_root=deck_dir)
        result.presentation.save(args.out)
    except RenderError as exc:
        print(f"error: {exc}", file=sys.stderr)
        return 2
    except OSError as exc:
        print(f"error: cannot write output {args.out}: {exc}", file=sys.stderr)
        return 2
    except (KeyError, ValueError, zipfile.BadZipFile) as exc:
        # structurally broken manifest or a corrupt template original
        print(f"error: invalid manifest/template structure: {exc!r}", file=sys.stderr)
        return 2

    s = result.stats
    print(f"rendered {s.pages} page(s) -> {args.out} "
          f"({s.formulas_native} native formula(s), {s.formulas_fallback} image fallback(s))")
    if s.fonts_embedded:
        fe = s.fonts_embedded
        parts = ", ".join(f"{name} {fe['weights'][name] // 1024}KB"
                          for name in efont.WEIGHT_ELEMENTS if name in fe["weights"])
        print(f"embedded {fe['face']} subset ({fe['chars']} chars): {parts}")
    elif s.font_warning:
        print(f"warning: {s.font_warning}", file=sys.stderr)
    return 0


if __name__ == "__main__":
    sys.exit(main())
