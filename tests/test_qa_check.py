"""Tests for the T6 QA gates: scripts/qa_check_pptx.py (fast, local) and
scripts/qa_check_web.py (e2e, needs npm + Playwright).

pptx-side contract under test: a correctly rendered fixture deck passes with
zero findings, while each injected defect class — placeholder residue,
over-budget text, page-count mismatch, title mismatch, off-canvas geometry,
crossed safe bottom — is caught by its own finding. The web-side contract:
a scaffolded project passes the Playwright checks (pages reachable, titles
aligned, KaTeX rendered, images loaded, no console errors), a broken image is
flagged, and --export-pptx produces a pptx that aligns with the web deck.
"""
import json
import shutil
import subprocess
import sys
from pathlib import Path

import pytest

REPO = Path(__file__).resolve().parent.parent
SCRIPTS_DIR = REPO / "scripts"
sys.path.insert(0, str(SCRIPTS_DIR))

import qa_check_pptx as qp  # noqa: E402
import render_pptx as rp  # noqa: E402

FIXTURE_DIR = REPO / "tests" / "fixtures"
FULL_DECK = FIXTURE_DIR / "deck.json"
PPTX_DECK = FIXTURE_DIR / "deck_pptx.json"
TEMPLATE_DIR = REPO / "templates" / "blcu-report"
TEMPLATE_PPTX = TEMPLATE_DIR / "blcu-report.pptx"
MANIFEST_PATH = TEMPLATE_DIR / "manifest.json"
RENDER_SCRIPT = SCRIPTS_DIR / "render_pptx.py"
QA_PPTX_SCRIPT = SCRIPTS_DIR / "qa_check_pptx.py"
QA_WEB_SCRIPT = SCRIPTS_DIR / "qa_check_web.py"
SCAFFOLD_SCRIPT = SCRIPTS_DIR / "scaffold_web.py"


def load(path):
    return json.loads(Path(path).read_text(encoding="utf-8"))


def absolutize_images(deck):
    """Point deck image paths at the fixture files (decks written outside
    fixtures/ would otherwise resolve images against their own directory)."""
    for page in deck["pages"]:
        for block in page["blocks"]:
            if block["type"] == "image":
                block["path"] = str((FIXTURE_DIR / block["path"]).resolve())
    return deck


def write_deck(deck, path):
    Path(path).write_text(json.dumps(deck, ensure_ascii=False), encoding="utf-8")
    return Path(path)


def render_cli(deck_path, out_path):
    return subprocess.run(
        [sys.executable, str(RENDER_SCRIPT), str(deck_path), "-o", str(out_path)],
        capture_output=True, text=True, encoding="utf-8", errors="replace",
    )


def qa_pptx(pptx, deck, *extra):
    return subprocess.run(
        [sys.executable, str(QA_PPTX_SCRIPT), str(pptx), "--deck", str(deck), *extra],
        capture_output=True, text=True, encoding="utf-8", errors="replace",
    )


xsl = rp.find_mml2omml_xsl()
requires_math = pytest.mark.skipif(xsl is None or not rp.HAS_LATEX2MATHML,
                                   reason="MML2OMML.XSL or latex2mathml unavailable")


# ---------------------------------------------------------------------------
# green path
# ---------------------------------------------------------------------------

@requires_math
def test_cli_green_on_rendered_fixture(tmp_path):
    out = tmp_path / "out.pptx"
    proc = render_cli(FULL_DECK, out)
    assert proc.returncode == 0, proc.stderr + proc.stdout  # not vacuously green
    proc = qa_pptx(out, FULL_DECK)
    assert proc.returncode == 0, proc.stdout
    assert "OK: 7 slide(s), 0 findings" in proc.stdout


@requires_math
def test_cli_green_json_format(tmp_path):
    out = tmp_path / "out.pptx"
    assert render_cli(FULL_DECK, out).returncode == 0
    proc = qa_pptx(out, FULL_DECK, "--format", "json")
    assert proc.returncode == 0
    report = json.loads(proc.stdout)
    assert report == {"valid": True, "slide_count": 7, "findings": []}


# ---------------------------------------------------------------------------
# injected defects are each caught by their own finding
# ---------------------------------------------------------------------------

@requires_math
def test_catches_placeholder_residue(tmp_path):
    deck = absolutize_images(load(FULL_DECK))
    deck["pages"][3]["blocks"][1]["text"] = "神经机器翻译在低资源场景下性能受限，TODO：此处填写数据。"
    deck_path = write_deck(deck, tmp_path / "residue.json")
    out = tmp_path / "residue.pptx"
    assert render_cli(deck_path, out).returncode == 0  # schema-legal: render passes
    proc = qa_pptx(out, deck_path)
    assert proc.returncode == 1
    findings = [line for line in proc.stdout.splitlines() if "[pptx.residue]" in line]
    assert any("[todo]" in line for line in findings)
    assert any("[此处填写]" in line for line in findings)
    assert "slides[4]" in findings[0]  # slide numbers are 1-based


def test_catches_over_budget_deck(tmp_path):
    # rendered through the API, bypassing the CLI's validation gate: the QA
    # revalidation is what must flag it（超限量从 manifest 读取，预算重校准后随动）
    deck = absolutize_images(load(FULL_DECK))
    manifest = load(MANIFEST_PATH)
    cap = manifest["archetypes"][deck["pages"][2]["archetype"]]["budget"]["text_total_max_chars"]
    deck["pages"][2]["blocks"].append({"type": "text", "text": "超" * (cap + 1)})
    result = rp.render_deck(deck, manifest, TEMPLATE_PPTX, image_root=FIXTURE_DIR)
    out = tmp_path / "over.pptx"
    result.presentation.save(out)
    proc = qa_pptx(out, write_deck(deck, tmp_path / "over.json"))
    assert proc.returncode == 1
    assert "[budget.text_block_chars]" in proc.stdout
    assert "[budget.text_total_chars]" in proc.stdout


@requires_math
def test_catches_page_count_mismatch(tmp_path):
    out = tmp_path / "three.pptx"
    assert render_cli(PPTX_DECK, out).returncode == 0  # 3 pages
    proc = qa_pptx(out, FULL_DECK)  # checked against the 7-page deck
    assert proc.returncode == 1
    assert "[pptx.page_count]" in proc.stdout


@requires_math
def test_catches_title_mismatch(tmp_path):
    deck = load(PPTX_DECK)
    out = tmp_path / "three.pptx"
    assert render_cli(PPTX_DECK, out).returncode == 0
    deck["pages"][1]["blocks"][0]["text"] = "被篡改的标题"
    proc = qa_pptx(out, write_deck(deck, tmp_path / "retitled.json"))
    assert proc.returncode == 1
    assert "[pptx.title_mismatch]" in proc.stdout
    assert "被篡改的标题" in proc.stdout


@requires_math
def test_catches_canvas_and_safe_bottom_overflows(tmp_path):
    from pptx import Presentation
    from pptx.util import Inches

    out = tmp_path / "base.pptx"
    assert render_cli(FULL_DECK, out).returncode == 0
    prs = Presentation(out)
    agenda_list = next(sp for sp in prs.slides[1].shapes if sp.name == "AgendaList")
    agenda_list.top = Inches(7.0)  # bottom lands at 12.2in, off-canvas
    text_area = next(sp for sp in prs.slides[2].shapes if sp.name == "ContentArea")
    text_area.top, text_area.height = Inches(4.32), Inches(3.0)  # bottom 7.32 > 6.9
    broken = tmp_path / "overflow.pptx"
    prs.save(broken)
    proc = qa_pptx(broken, FULL_DECK)
    assert proc.returncode == 1
    assert "[pptx.canvas_overflow]" in proc.stdout
    assert "[pptx.content_bottom]" in proc.stdout


# ---------------------------------------------------------------------------
# style layer (#20): font consistency / weight hierarchy / color legality /
# brand layer, each caught by its own finding on a post-render mutation
# ---------------------------------------------------------------------------

def _mutated_fixture(tmp_path, name="base.pptx"):
    """Render FULL_DECK and open the result for run-level mutation."""
    from pptx import Presentation

    out = tmp_path / name
    assert render_cli(FULL_DECK, out).returncode == 0
    return out, Presentation(out)


def _shape_runs(prs, slide_i, shape_name):
    sp = next(sp for sp in prs.slides[slide_i].shapes if sp.name == shape_name)
    return sp, sp._element.findall(f".//{{{qp.NS_A}}}r")


@requires_math
def test_catches_font_face_drift(tmp_path):
    # a run restyled to another face after rendering: page-to-page font
    # consistency is a gate, not a preference
    out, prs = _mutated_fixture(tmp_path)
    _, runs = _shape_runs(prs, 2, "ContentArea")
    rPr = runs[0].find(f"{{{qp.NS_A}}}rPr")
    for tag in ("latin", "ea"):
        rPr.find(f"{{{qp.NS_A}}}{tag}").set("typeface", "微软雅黑")
    broken = tmp_path / "face.pptx"
    prs.save(broken)
    proc = qa_pptx(broken, FULL_DECK)
    assert proc.returncode == 1
    findings = [l for l in proc.stdout.splitlines() if "[pptx.font_face]" in l]
    assert len(findings) == 2  # latin + ea both drifted
    assert "微软雅黑" in findings[0]


@requires_math
def test_catches_weight_hierarchy_break(tmp_path):
    # title demoted to regular, and a body run bolded without the emphasis
    # color: both break the weight hierarchy the tokens declare
    out, prs = _mutated_fixture(tmp_path)
    title = next(sp for sp in prs.slides[2].shapes if qp.ph_type(sp) == "title")
    title._element.find(f".//{{{qp.NS_A}}}rPr").set("b", "0")
    _, runs = _shape_runs(prs, 2, "ContentArea")
    runs[0].find(f"{{{qp.NS_A}}}rPr").set("b", "1")  # stays black: not emphasis
    broken = tmp_path / "weights.pptx"
    prs.save(broken)
    proc = qa_pptx(broken, FULL_DECK)
    assert proc.returncode == 1
    findings = [l for l in proc.stdout.splitlines() if "[pptx.role_style]" in l]
    assert any("role 'title'" in l for l in findings)
    assert any("font-weight" in l and "ContentArea" in l for l in findings)


@requires_math
def test_catches_illegal_text_colors(tmp_path):
    # band green as a text color is always illegal; accent on a regular run
    # is legal only as the bold emphasis variant
    out, prs = _mutated_fixture(tmp_path)
    _, runs = _shape_runs(prs, 4, "TextArea")  # text-image text region: 3 runs
    runs[0].find(f"{{{qp.NS_A}}}rPr/{{{qp.NS_A}}}solidFill/{{{qp.NS_A}}}srgbClr") \
        .set("val", "C5E0B4")
    runs[1].find(f"{{{qp.NS_A}}}rPr/{{{qp.NS_A}}}solidFill/{{{qp.NS_A}}}srgbClr") \
        .set("val", "548235")  # b stays 0: accent without bold is a misuse
    broken = tmp_path / "colors.pptx"
    prs.save(broken)
    proc = qa_pptx(broken, FULL_DECK)
    assert proc.returncode == 1
    findings = [l for l in proc.stdout.splitlines() if "[pptx.text_color]" in l]
    assert findings and "C5E0B4" in findings[0]  # the band color is never text
    style = [l for l in proc.stdout.splitlines() if "[pptx.role_style]" in l]
    assert any("#548235" in l for l in style)  # accent on a regular run flagged


@requires_math
def test_catches_missing_brand_layer(tmp_path):
    # the content master's corner logo removed after rendering: the slide
    # inherits a master without the measured brand layer
    from pptx.util import Emu

    out, prs = _mutated_fixture(tmp_path)
    master = prs.slides[2].slide_layout.slide_master
    logo = next(
        sp for sp in master.shapes
        if sp.shape_type == qp.PICTURE_SHAPE_TYPE
        and abs(Emu(sp.left).inches - 0.27) < 0.05
        and abs(Emu(sp.top).inches - 6.42) < 0.05)
    logo._element.getparent().remove(logo._element)
    broken = tmp_path / "brandless.pptx"
    prs.save(broken)
    proc = qa_pptx(broken, FULL_DECK)
    assert proc.returncode == 1
    findings = [l for l in proc.stdout.splitlines() if "[pptx.brand_layer]" in l]
    assert findings and "image7.png" in findings[0]


@requires_math
def test_title_face_variant_is_an_explicit_machine_checked_override(tmp_path):
    # the 思源黑 vs 思源宋 A/B flow: a variant manifest overriding the title
    # role's face renders serif titles that pass the variant's own gate, and
    # the same render against the committed manifest is font-face drift
    variant = load(MANIFEST_PATH)
    for name in ("title", "title_long"):
        variant["typography"]["tokens"]["roles"][name]["face"] = "Noto Serif SC"
    variant_path = tmp_path / "manifest-serif.json"
    variant_path.write_text(json.dumps(variant, ensure_ascii=False), encoding="utf-8")

    deck = absolutize_images(load(FULL_DECK))
    deck_path = write_deck(deck, tmp_path / "deck.json")
    result = rp.render_deck(deck, variant, TEMPLATE_PPTX, image_root=FIXTURE_DIR)
    out = tmp_path / "serif.pptx"
    result.presentation.save(out)

    proc = qa_pptx(out, deck_path, "--manifest", variant_path)
    assert proc.returncode == 0, proc.stdout  # the override is machine-checked
    proc = qa_pptx(out, deck_path)
    assert proc.returncode == 1  # committed tokens still say Noto Sans SC
    assert "[pptx.font_face]" in proc.stdout
    assert "Noto Serif SC" in proc.stdout


@requires_math
def test_catches_top_anchored_pure_text_page(tmp_path):
    # the pure-text design centers pages without formulas; a box flipped
    # back to top anchor after rendering reintroduces the bottom void
    out, prs = _mutated_fixture(tmp_path)
    box = next(sp for sp in prs.slides[3].shapes if sp.name == "ContentArea")
    bodyPr = box._element.find(
        f"{{{qp.NS_P}}}txBody/{{{qp.NS_A}}}bodyPr")
    bodyPr.set("anchor", "t")
    broken = tmp_path / "topanchor.pptx"
    prs.save(broken)
    proc = qa_pptx(broken, FULL_DECK)
    assert proc.returncode == 1
    findings = [l for l in proc.stdout.splitlines() if "[pptx.role_style]" in l]
    assert any("ContentArea" in l and "centers" in l for l in findings)


@requires_math
def test_style_checks_stay_total_on_malformed_manifest(tmp_path):
    # hand-edited variant manifests are a supported flow; a malformed
    # tokens/brand_layer section must degrade to findings-or-skip, never a
    # traceback (the same contract the deck-side checks already hold)
    out = tmp_path / "out.pptx"
    assert render_cli(FULL_DECK, out).returncode == 0
    m = load(MANIFEST_PATH)
    m["typography"]["tokens"]["roles"]["title"] = "44pt"      # role not a dict
    m["typography"]["tokens"]["emphasis"] = ["**"]            # not a dict
    m["typography"]["tokens"]["role_bindings"] = {            # leaf not a string
        "text-formula": {"text": {"role": "body"}, "title": ["title"]}}
    m["typography"]["tokens"]["roles"]["title_long"]["over_chars"] = "15"  # not numeric
    m["brand_layer"] = {"content": "x", "cover": {"mid_band": "band", "logos": 3}}
    mp = tmp_path / "malformed-manifest.json"
    mp.write_text(json.dumps(m, ensure_ascii=False), encoding="utf-8")
    proc = qa_pptx(out, FULL_DECK, "--manifest", mp)
    assert "Traceback" not in proc.stderr
    assert proc.returncode in (0, 1)


def test_exit_2_on_usage_and_io_errors(tmp_path):
    proc = qa_pptx(tmp_path / "ghost.pptx", FULL_DECK)
    assert proc.returncode == 2
    assert "pptx not found" in proc.stderr
    proc = subprocess.run(
        [sys.executable, str(QA_PPTX_SCRIPT), str(tmp_path / "any.pptx"),
         "--deck", str(tmp_path / "no-deck.json")],
        capture_output=True, text=True, encoding="utf-8", errors="replace",
    )
    assert proc.returncode == 2
    assert "Traceback" not in proc.stderr


# ---------------------------------------------------------------------------
# unit seam: the residue matcher
# ---------------------------------------------------------------------------

def test_find_residue_matches_fillers_without_false_positives():
    for text in ("TODO", "todo:", "FIXME now", "meeting TBD", "xxxx", "lorem ipsum",
                 "[insert name]", "<insert>here", "a placeholder b", "Sample Text",
                 "待补充", "此处填写内容", "示例文本"):
        assert qp.find_residue(text), text
    for clean in ("拓扑排序", "contrastive learning", "TODOor 的边界情况",
                  "maximization", "EXP（指数）", "插入 insert 图片", "西安 TBDX 公司",
                  "可学习占位符", "占位效应", "待定系数法", "时间待定"):
        # "TODOor" / "TBDX" / bare "insert" stay legal: the patterns need word
        # boundaries or a bracket tag; bare 占位/待定 stay legal because they
        # live in real academic terms (可学习占位符 / 待定系数法)
        assert not qp.find_residue(clean), clean


def test_qa_pptx_stays_total_on_schema_broken_deck(tmp_path):
    # a deck edited after rendering (page missing its blocks) must yield a
    # findings report, never a traceback — the tool exists to judge such decks
    out = tmp_path / "out.pptx"
    assert render_cli(PPTX_DECK, out).returncode == 0
    deck = {"template": "blcu-report", "pages": [{"archetype": "cover"}]}
    proc = qa_pptx(out, write_deck(deck, tmp_path / "broken.json"))
    assert proc.returncode == 1
    assert "Traceback" not in proc.stderr
    assert "schema" in proc.stdout


def test_web_qa_reports_invalid_deck_without_a_browser(tmp_path):
    # schema-broken project decks are caught by validation before any
    # dev server / browser starts
    web = tmp_path / "web"
    (web / "src").mkdir(parents=True)
    (web / "package.json").write_text("{}", encoding="utf-8")
    (web / "src" / "deck.json").write_text(
        json.dumps({"template": "blcu-report", "pages": [{"archetype": "cover"}]}),
        encoding="utf-8")
    shutil.copyfile(MANIFEST_PATH, web / "src" / "manifest.json")
    (web / "node_modules").mkdir()  # present, so the earlier env check passes
    proc = subprocess.run(
        [sys.executable, str(QA_WEB_SCRIPT), str(web)],
        capture_output=True, text=True, encoding="utf-8", errors="replace",
    )
    assert proc.returncode == 1
    assert "schema.missing_field" in proc.stdout
    assert "Traceback" not in proc.stderr


# ---------------------------------------------------------------------------
# e2e web QA (needs npm for the dev server + playwright with Chromium)
# ---------------------------------------------------------------------------

npm = shutil.which("npm")
try:
    import playwright.sync_api as _pw  # noqa: F401

    has_playwright = True
except ImportError:
    has_playwright = False
requires_web = pytest.mark.skipif(
    npm is None or not has_playwright,
    reason="npm or python playwright unavailable")


@requires_web
def test_web_qa_green_then_defect_then_export(tmp_path):
    web = tmp_path / "web"
    proc = subprocess.run(
        [sys.executable, str(SCAFFOLD_SCRIPT), str(FULL_DECK), "-o", str(web)],
        capture_output=True, text=True, encoding="utf-8", errors="replace",
    )
    assert proc.returncode == 0, proc.stderr + proc.stdout
    install = subprocess.run(
        [npm, "install", "--no-audit", "--no-fund", "--loglevel=error"],
        cwd=web, capture_output=True, text=True, encoding="utf-8",
        errors="replace", timeout=600,
    )
    assert install.returncode == 0, install.stderr + install.stdout

    shots = tmp_path / "shots"
    exported = tmp_path / "exported.pptx"
    proc = subprocess.run(
        [sys.executable, str(QA_WEB_SCRIPT), str(web),
         "--screenshots", str(shots), "--export-pptx", str(exported)],
        capture_output=True, text=True, encoding="utf-8", errors="replace",
        timeout=300,
    )
    assert proc.returncode == 0, proc.stdout + proc.stderr
    assert "0 findings" in proc.stdout
    assert exported.is_file()
    assert sorted(p.name for p in shots.glob("page-*.png")) == [
        f"page-{i:02d}.png" for i in range(1, 8)]

    # defect phase: a missing deck image must fail the gate (caught by the
    # upfront deck validation, before any browser is involved)
    (web / "public" / "material" / "images" / "main-results-bar.png").unlink()
    proc = subprocess.run(
        [sys.executable, str(QA_WEB_SCRIPT), str(web)],
        capture_output=True, text=True, encoding="utf-8", errors="replace",
        timeout=300,
    )
    assert proc.returncode == 1
    assert "image.path_missing" in proc.stdout

    # defect phase 2: a missing brand asset must fail the brand-layer gate
    shutil.copyfile(FIXTURE_DIR / "material" / "images" / "main-results-bar.png",
                    web / "public" / "material" / "images" / "main-results-bar.png")
    (web / "public" / "brand" / "image7.png").unlink()
    proc = subprocess.run(
        [sys.executable, str(QA_WEB_SCRIPT), str(web)],
        capture_output=True, text=True, encoding="utf-8", errors="replace",
        timeout=300,
    )
    assert proc.returncode == 1
    assert "web.brand_asset_broken" in proc.stdout

    # defect phase 3: a hard-coded style override bypassing the tokens must
    # fail the style gate (title weight + emphasis color)
    shutil.copyfile(REPO / "assets" / "web-template" / "public" / "brand" / "image7.png",
                    web / "public" / "brand" / "image7.png")
    override = ('\n[data-role="title"] { font-weight: 400 !important; }\n'
                '\n.emph { color: #ff0000 !important; }\n')
    css = web / "src" / "styles.css"
    css.write_text(css.read_text(encoding="utf-8") + override, encoding="utf-8")
    proc = subprocess.run(
        [sys.executable, str(QA_WEB_SCRIPT), str(web)],
        capture_output=True, text=True, encoding="utf-8", errors="replace",
        timeout=300,
    )
    assert proc.returncode == 1
    assert "web.style" in proc.stdout
    css.write_text(css.read_text(encoding="utf-8").replace(override, ""),
                   encoding="utf-8")

    # a white-background screenshot carries the token hairline and stays
    # green: renderer detection, outline and gate all exercised on it
    shutil.copyfile(FIXTURE_DIR / "material" / "images" / "white-chart.png",
                    web / "public" / "material" / "images" / "main-results-bar.png")
    proc = subprocess.run(
        [sys.executable, str(QA_WEB_SCRIPT), str(web)],
        capture_output=True, text=True, encoding="utf-8", errors="replace",
        timeout=300,
    )
    assert proc.returncode == 0, proc.stdout + proc.stderr

    # defect phase 4: a hand-corrupted manifest copy (the web sibling of the
    # pptx gate's malformed-manifest contract) must degrade to findings,
    # never a traceback. The mutations stay JS-benign — the React app still
    # mounts (role_bindings container malformations crash layout.js at mount,
    # so those are covered by the unit tests below instead).
    m = json.loads((web / "src" / "manifest.json").read_text(encoding="utf-8"))
    m["typography"]["tokens"]["roles"]["title"] = "44pt"
    m["typography"]["tokens"]["roles"]["title_long"]["over_chars"] = "15"
    m["typography"]["tokens"]["roles"]["body"]["color"] = "red"  # named CSS color
    m["typography"]["tokens"]["emphasis"] = ["**"]
    m["brand_layer"] = {"content": {"top_band": "band", "corner_logo": "x"},
                        "cover": {"mid_band": "band", "logos": [],
                                  "corner_logo_bar": "bar"}}
    (web / "src" / "manifest.json").write_text(
        json.dumps(m, ensure_ascii=False), encoding="utf-8")
    proc = subprocess.run(
        [sys.executable, str(QA_WEB_SCRIPT), str(web)],
        capture_output=True, text=True, encoding="utf-8", errors="replace",
        timeout=300,
    )
    assert "Traceback" not in proc.stderr + proc.stdout
    assert proc.returncode in (0, 1)
    # the same malformed manifest through the export chain: the renderer is
    # strict on manifest structure, so the chain fails — cleanly, exit 2
    proc = subprocess.run(
        [sys.executable, str(QA_WEB_SCRIPT), str(web),
         "--export-pptx", str(tmp_path / "corrupt.pptx")],
        capture_output=True, text=True, encoding="utf-8", errors="replace",
        timeout=300,
    )
    assert proc.returncode == 2
    assert "export chain failed" in proc.stderr
    assert "Traceback" not in proc.stderr


def test_hex_rgb_parses_only_six_hex():
    import qa_check_web as qw

    assert qw._hex_rgb("#548235") == (0x54, 0x82, 0x35)
    assert qw._hex_rgb("548235") == (0x54, 0x82, 0x35)
    # anything else is no-expectation, never a raise: named colors, quoted
    # junk, non-strings
    for bad in ("red", "44pt", "", "#12345", None, 44, ["548235"]):
        assert qw._hex_rgb(bad) is None, bad


def test_expected_title_role_total_on_malformed_containers():
    # role_bindings/roles container malformations crash layout.js at mount,
    # so the browser e2e cannot carry them — the unit seam holds the contract
    import qa_check_web as qw

    page = {"archetype": "text-formula"}
    for tokens in ({"role_bindings": "cover", "roles": {}},
                   {"role_bindings": ["x"], "roles": "roles"},
                   {"role_bindings": {"text-formula": "x"}},
                   {"roles": {"title_long": "x"}}):
        assert qw._expected_title_role(tokens, page, "任何标题") in (None, "title")
    # the downgrade still fires through numeric-string coercion (layout.js
    # coerces too — all three readers agree)
    tokens = {"role_bindings": {"text-formula": {"title": "title"}},
              "roles": {"title_long": {"over_chars": "15"}}}
    assert qw._expected_title_role(
        tokens, page, "这是一个超过十五个字符的超长标题演示") == "title_long"


def test_web_qa_exit_2_without_node_modules(tmp_path):
    web = tmp_path / "bare"
    (web / "src").mkdir(parents=True)
    (web / "package.json").write_text("{}", encoding="utf-8")
    (web / "src" / "deck.json").write_text('{"template": "blcu-report", "pages": []}',
                                           encoding="utf-8")
    (web / "src" / "manifest.json").write_text("{}", encoding="utf-8")
    proc = subprocess.run(
        [sys.executable, str(QA_WEB_SCRIPT), str(web)],
        capture_output=True, text=True, encoding="utf-8", errors="replace",
    )
    assert proc.returncode == 2
    assert "npm install" in proc.stderr
