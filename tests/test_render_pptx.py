"""Seam-1 pipeline tests for scripts/render_pptx.py (all 6 archetypes).

Contract under test: a deck.json renders against the real template original via
Clone & Fill; the output reopens with python-pptx carrying the right page
count, per-page titles, native OMML formulas, embedded images, manifest-region
geometry, master branding, zero template residue — and, since S3, the full
typography.tokens role system at run level: face/size/weight/color per role,
**keyword** emphasis runs, caption scrim strips, the body spacing rhythm
(line 1.5 / before 12pt), token textbox insets, hairline on white-background
images, and the S1 font embed riding along. deck_pptx.json exercises cover +
text-formula in depth (T3); deck.json is the T2 full-archetype fixture covering
agenda / text-image / chart-focus / closing too (T4). The template original
must be byte-identical after rendering.
"""
import hashlib
import json
import subprocess
import sys
import zipfile
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Inches

REPO = Path(__file__).resolve().parent.parent
SCRIPTS_DIR = REPO / "scripts"
sys.path.insert(0, str(SCRIPTS_DIR))

import render_pptx as rp  # noqa: E402

FIXTURE_DECK = REPO / "tests" / "fixtures" / "deck_pptx.json"
FULL_DECK = REPO / "tests" / "fixtures" / "deck.json"
FIXTURE_DIR = REPO / "tests" / "fixtures"
TEMPLATE_DIR = REPO / "templates" / "blcu-report"
TEMPLATE_PPTX = TEMPLATE_DIR / "blcu-report.pptx"
MANIFEST_PATH = TEMPLATE_DIR / "manifest.json"
SCRIPT = SCRIPTS_DIR / "render_pptx.py"

NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
NS_M = "http://schemas.openxmlformats.org/officeDocument/2006/math"
NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"

# Original text from the template's clone-source slides (1, 2, 3, 7, 18):
# none of it may survive into a rendered deck.
TEMPLATE_RESIDUE = [
    # slide 1 (cover / closing source)
    "阿里天池亚军方案复现",
    "汇报人:赵法科",
    # slide 2 (agenda source)
    "算法思想",
    "一 预处理：预处理",
    "序列×参数并行评估",
    # slide 3 (text-formula source)
    "预处理（简化",
    "旋转枚举）",
    "原版做法",
    "简化：严格外包简化",
    "凹顶点删除",
    "凸顶点对合并",
    # slide 7 (text-image source)
    "预处理（NFP+IFP+NFP预计算）",
    "待放件",
    "Minko",
    # slide 18 (chart-focus source)
    "实验结果+问题分析",
    "和autonest对比",
    "packingUI",
]


def load(name):
    return json.loads((REPO / "tests" / "fixtures" / name).read_text(encoding="utf-8"))


def render(tmp_path, deck_path=None):
    """Run the CLI on the fixture deck; return (out_path, subprocess result)."""
    deck_path = Path(deck_path or FIXTURE_DECK)
    out_path = tmp_path / "out.pptx"
    proc = subprocess.run(
        [sys.executable, str(SCRIPT), str(deck_path), "-o", str(out_path)],
        capture_output=True, text=True, encoding="utf-8", errors="replace",
    )
    return out_path, proc


def omath_count(slide):
    return len(slide.part._element.findall(f".//{{{NS_M}}}oMath"))


def ph_type(sp):
    ph = sp._element.find(f".//{{{NS_P}}}ph")
    return None if ph is None else (ph.get("type") or "body")


def shapes(slide):
    return list(slide.shapes)


def by_ph(slide, want):
    for sp in shapes(slide):
        if ph_type(sp) == want:
            return sp
    return None


def non_placeholder_shapes(slide):
    return [sp for sp in shapes(slide) if ph_type(sp) is None]


def all_text(slide):
    """Text of every non-sldNum shape, a:t and m:t alike."""
    parts = []
    for sp in shapes(slide):
        if ph_type(sp) == "sldNum":
            continue
        parts.extend(t.text or "" for t in sp._element.iter(f"{{{NS_A}}}t"))
        parts.extend(t.text or "" for t in sp._element.iter(f"{{{NS_M}}}t"))
    return "".join(parts)


def region_of(sp):
    return (sp.left, sp.top, sp.width, sp.height)


def assert_region(sp, region):
    assert list(region_of(sp)) == pytest.approx(
        [Inches(region[k]) for k in ("x", "y", "w", "h")], abs=2
    )


def assert_fitted_in(pic, slot):
    """Picture contain-fits its slot: inside, centered, tight on one axis."""
    sx, sy, sw, sh = (Inches(slot[k]) for k in ("x", "y", "w", "h"))
    left, top, width, height = region_of(pic)
    assert left >= sx - 2 and top >= sy - 2
    assert left + width <= sx + sw + 2 and top + height <= sy + sh + 2
    # contain-fit is tight on the binding axis
    assert width == pytest.approx(sw, abs=2) or height == pytest.approx(sh, abs=2)
    # centered on both axes
    assert left - sx == pytest.approx(sx + sw - (left + width), abs=4)
    assert top - sy == pytest.approx(sy + sh - (top + height), abs=4)


def run_color(rPr):
    fill = rPr.find(f"{{{NS_A}}}solidFill/{{{NS_A}}}srgbClr")
    return fill.get("val") if fill is not None else None


def run_face(rPr, which):
    el = rPr.find(f"{{{NS_A}}}{which}")
    return el.get("typeface") if el is not None else None


def run_style(sp):
    """(sz_pt, bold, color, latin, ea) of the first explicitly sized run."""
    for rPr in sp._element.iter(f"{{{NS_A}}}rPr"):
        sz = rPr.get("sz")
        if sz is None:
            continue
        return (int(sz) // 100, rPr.get("b") == "1", run_color(rPr),
                run_face(rPr, "latin"), run_face(rPr, "ea"))
    return None


def shape_runs(sp):
    """[(text, sz_pt, bold, color)] for every text run in the shape, in order."""
    out = []
    for r in sp._element.iter(f"{{{NS_A}}}r"):
        rPr = r.find(f"{{{NS_A}}}rPr")
        out.append((r.find(f"{{{NS_A}}}t").text,
                    int(rPr.get("sz", "0")) // 100,
                    rPr.get("b") == "1",
                    run_color(rPr)))
    return out


def first_para_props(sp):
    """(line_spacing_pct, space_before_pt) of the shape's first paragraph."""
    pPr = sp._element.find(f".//{{{NS_A}}}pPr")
    if pPr is None:
        return (None, 0)
    ln = pPr.find(f"{{{NS_A}}}lnSpc/{{{NS_A}}}spcPct")
    bef = pPr.find(f"{{{NS_A}}}spcBef/{{{NS_A}}}spcPts")
    return (int(ln.get("val")) // 1000 if ln is not None else None,
            int(bef.get("val")) // 100 if bef is not None else 0)


def prose_para_props(sp):
    """(line_spacing_pct, space_before_pt) of the first a:r-carrying paragraph
    (math paragraphs lead formula-first content flows)."""
    for p_el in sp._element.iter(f"{{{NS_A}}}p"):
        if p_el.find(f"{{{NS_A}}}r") is None:
            continue
        pPr = p_el.find(f"{{{NS_A}}}pPr")
        ln = pPr.find(f"{{{NS_A}}}lnSpc/{{{NS_A}}}spcPct") if pPr is not None else None
        bef = pPr.find(f"{{{NS_A}}}spcBef/{{{NS_A}}}spcPts") if pPr is not None else None
        return (int(ln.get("val")) // 1000 if ln is not None else None,
                int(bef.get("val")) // 100 if bef is not None else 0)
    return (None, 0)


def body_insets(sp):
    bodyPr = sp._element.find(f".//{{{NS_A}}}bodyPr")
    return tuple(int(bodyPr.get(k)) for k in ("lIns", "tIns", "rIns", "bIns"))


NOTO = "Noto Sans SC"
TITLE_44 = (44, True, "000000", NOTO, NOTO)
TITLE_28 = (28, True, "000000", NOTO, NOTO)


xsl = rp.find_mml2omml_xsl()
requires_math = pytest.mark.skipif(xsl is None or not rp.HAS_LATEX2MATHML,
                                  reason="MML2OMML.XSL or latex2mathml unavailable")


# ---------------------------------------------------------------------------
# happy path: fixture deck renders, reopens, and every assertion holds
# ---------------------------------------------------------------------------

@requires_math
def test_cli_renders_fixture_deck(tmp_path):
    out_path, proc = render(tmp_path)
    assert proc.returncode == 0, proc.stderr + proc.stdout
    assert out_path.is_file()


@requires_math
def test_render_pipeline(tmp_path):
    deck = load("deck_pptx.json")
    manifest = json.loads(MANIFEST_PATH.read_text(encoding="utf-8"))
    result = rp.render_deck(deck, manifest, TEMPLATE_PPTX)
    prs = result.presentation
    arch = manifest["archetypes"]

    # -- page count and order
    assert len(prs.slides) == len(deck["pages"]) == 3

    # -- branding survives: layout names + all three masters + master media
    assert [s.slide_layout.name for s in prs.slides] == ["三logo标题页", "标题和内容", "标题和内容"]
    assert len(prs.slide_masters) == 3
    out_path = tmp_path / "pipeline.pptx"
    prs.save(out_path)
    with zipfile.ZipFile(out_path) as z:
        names = z.namelist()
        slides = sorted(n for n in names if n.startswith("ppt/slides/slide"))
        assert len(slides) == 3
        assert len([n for n in names if n.startswith("ppt/slideMasters/slideMaster")]) == 3
        # brand media from the masters stays
        for img in ("image1.png", "image2.gif", "image3.png"):
            assert f"ppt/media/{img}" in names
        # cloned slides carry no speaker notes (演讲稿 is a separate deliverable)
        assert not [n for n in names if n.startswith("ppt/notesSlides/")]

    # -- slide 1: cover
    s1 = prs.slides[0]
    title1 = by_ph(s1, "title")
    assert title1 is not None
    assert title1.text_frame.text == deck["pages"][0]["blocks"][0]["text"]
    assert_region(title1, arch["cover"]["regions"]["title"])
    assert run_style(title1) == TITLE_44
    sub1 = by_ph(s1, "body")
    assert sub1 is not None
    assert sub1.text_frame.text == "汇报人:张三 2026-08-25"
    assert_region(sub1, arch["cover"]["regions"]["subtitle"])
    assert run_style(sub1) == (24, False, "000000", NOTO, NOTO)
    assert by_ph(s1, "sldNum") is None  # cover carries no page number
    assert non_placeholder_shapes(s1) == []

    # -- slide 2: text-formula with 2 native OMML formulas + text
    s2 = prs.slides[1]
    title2 = by_ph(s2, "title")
    assert title2.text_frame.text == "方法：对比学习框架"
    assert_region(title2, arch["text-formula"]["regions"]["title"])
    assert run_style(title2) == TITLE_44
    assert omath_count(s2) == 2
    # page-number placeholder preserved as a field
    sldnum = by_ph(s2, "sldNum")
    assert sldnum is not None
    assert sldnum._element.find(f".//{{{NS_A}}}fld") is not None
    # injected math runs: Cambria Math at the manifest formula size
    math_rprs = [rPr for r in s2.part._element.iter(f"{{{NS_M}}}r")
                 for rPr in r.findall(f"{{{NS_A}}}rPr")]
    assert math_rprs
    for rPr in math_rprs:
        assert rPr.get("sz") == "2000"
        latin = rPr.find(f"{{{NS_A}}}latin")
        assert latin is not None and latin.get("typeface") == "Cambria Math"
    # content: ONE full-height flow box — formulas and prose interleaved
    # (the template author's idiom; S6 density redesign)
    boxes = non_placeholder_shapes(s2)
    assert len(boxes) == 1
    flow = boxes[0]
    assert flow.name == "ContentArea"
    assert_region(flow, arch["text-formula"]["regions"]["content"])
    # the prose paragraph carries the body style (math runs style themselves)
    paras = flow.text_frame.paragraphs
    assert paras[-1].text == deck["pages"][1]["blocks"][3]["text"]
    prose_rPr = paras[-1]._p.find(f"{{{NS_A}}}r/{{{NS_A}}}rPr")
    assert (int(prose_rPr.get("sz")) // 100, prose_rPr.get("b")) == (20, "0")
    ea = prose_rPr.find(f"{{{NS_A}}}ea")
    assert ea is not None and ea.get("typeface") == NOTO
    # the content flow vertically centers in the full-height region
    anchor = flow._element.find(
        f"{{{NS_P}}}txBody/{{{NS_A}}}bodyPr").get("anchor")
    assert anchor == "ctr"

    # -- slide 3: pure-text variant centers in the full-height flow region
    s3 = prs.slides[2]
    assert omath_count(s3) == 0
    assert by_ph(s3, "sldNum") is not None
    boxes3 = non_placeholder_shapes(s3)
    assert len(boxes3) == 1
    assert_region(boxes3[0], arch["text-formula"]["regions"]["content"])
    anchor = boxes3[0]._element.find(
        f"{{{NS_P}}}txBody/{{{NS_A}}}bodyPr").get("anchor")
    assert anchor == "ctr"
    assert run_style(boxes3[0]) == (20, False, "000000", NOTO, NOTO)
    # markers are stripped from the rendered text (they become styled runs)
    assert boxes3[0].text_frame.text == (
        "神经机器翻译在低资源场景下性能受限，平行语料稀缺是主要瓶颈。\n"
        "对比学习通过拉近正例、推远负例，可在无监督条件下学到通用表示。"
    )

    # -- slide size untouched (EMU truth: 13.3333 x 7.5 in; manifest 13.33 is rounded)
    assert (prs.slide_width, prs.slide_height) == (12192000, 6858000)

    # -- zero template residue anywhere in the deck
    for slide in prs.slides:
        text = all_text(slide)
        for residue in TEMPLATE_RESIDUE:
            assert residue not in text


@requires_math
def test_template_original_is_never_modified(tmp_path):
    before = hashlib.sha256(TEMPLATE_PPTX.read_bytes()).hexdigest()
    _, proc = render(tmp_path)
    assert proc.returncode == 0, proc.stderr + proc.stdout  # not vacuously green
    after = hashlib.sha256(TEMPLATE_PPTX.read_bytes()).hexdigest()
    assert before == after


@requires_math
def test_long_title_switches_to_28pt(tmp_path):
    deck = load("deck_pptx.json")
    # 16 CJK chars: over long_title_over_chars (15), within title_max_chars (25)
    deck["pages"][2]["blocks"][0]["text"] = "十六字标题测试用例一二三四五六七八"
    long_deck = tmp_path / "long.json"
    long_deck.write_text(json.dumps(deck, ensure_ascii=False), encoding="utf-8")
    out_path, proc = render(tmp_path, long_deck)
    assert proc.returncode == 0, proc.stderr + proc.stdout
    prs = Presentation(out_path)
    assert run_style(by_ph(prs.slides[2], "title")) == TITLE_28


# ---------------------------------------------------------------------------
# formula chain
# ---------------------------------------------------------------------------

@requires_math
def test_latex_to_omml_produces_native_math():
    omml = rp.convert_latex_to_omml(r"\mathcal{L} = \frac{a}{b}")
    assert omml is not None
    assert omml.tag == f"{{{NS_M}}}oMath"
    assert omml.findall(f".//{{{NS_M}}}r")


@requires_math
def test_norm_bars_keep_their_operands_through_the_chain():
    # MML2OMML.XSL reads adjacent ‖…‖ pairs as one delimiter run and drops
    # the enclosed operands (‖u‖ ‖v‖ rendered as empty bars), and drops
    # <mspace> outright (gluing the two groups together); the chain
    # normalizes bare \lVert/\rVert to \left/\right and \, to ~ so both the
    # operands and the separating space survive. Found comparing web vs
    # pptx renders side by side (#19 verification).
    latex = r"\frac{u^{\top} v}{\lVert u \rVert \, \lVert v \rVert}"
    omml = rp.convert_latex_to_omml(latex)
    texts = [t.text or "" for t in omml.iter(f"{{{NS_M}}}t")]
    assert texts.count("u") == 2 and texts.count("v") == 2  # numerator + denominator
    assert "\xa0" in texts  # the thin space survives as a real space run


def test_formula_fallback_renders_images(tmp_path, monkeypatch):
    monkeypatch.setattr(rp, "convert_latex_to_omml", lambda latex, xsl_path=None: None)
    deck = load("deck_pptx.json")
    manifest = json.loads(MANIFEST_PATH.read_text(encoding="utf-8"))
    result = rp.render_deck(deck, manifest, TEMPLATE_PPTX)

    s2 = result.presentation.slides[1]
    assert omath_count(s2) == 0
    content = manifest["archetypes"]["text-formula"]["regions"]["content"]
    pics = [sp for sp in non_placeholder_shapes(s2) if sp.shape_type == 13]  # PICTURE
    assert len(pics) == 2
    assert result.stats.formulas_fallback == 2
    for pic in pics:
        # fallback images stack inside the content region's lower band,
        # clear of the native-paragraph headroom above
        assert pic.top >= Inches(content["y"]) - 2
        assert pic.top + pic.height <= Inches(content["y"]) + Inches(content["h"]) + 2


# ---------------------------------------------------------------------------
# gate behavior: invalid decks and out-of-scope archetypes are refused
# ---------------------------------------------------------------------------

def test_cli_refuses_invalid_deck(tmp_path):
    deck = load("deck_pptx.json")
    # 26 CJK chars > text-formula title_max_chars (25)
    deck["pages"][1]["blocks"][0]["text"] = "超" * 26
    bad = tmp_path / "bad.json"
    bad.write_text(json.dumps(deck, ensure_ascii=False), encoding="utf-8")
    out_path, proc = render(tmp_path, bad)
    assert proc.returncode == 1
    assert "budget" in proc.stdout
    assert not out_path.exists()


def test_render_deck_refuses_unimplemented_archetype():
    deck = load("deck.json")
    # a validation-clean shape: only render_deck's archetype scope can refuse it
    deck["pages"][0]["archetype"] = "unicorn"
    manifest = json.loads(MANIFEST_PATH.read_text(encoding="utf-8"))
    with pytest.raises(rp.RenderError) as excinfo:
        rp.render_deck(deck, manifest, TEMPLATE_PPTX, image_root=FIXTURE_DIR)
    message = str(excinfo.value)
    assert "unicorn" in message
    assert "cover" in message and "closing" in message


def test_cli_refuses_to_overwrite_template(tmp_path):
    before = hashlib.sha256(TEMPLATE_PPTX.read_bytes()).hexdigest()
    proc = subprocess.run(
        [sys.executable, str(SCRIPT), str(FIXTURE_DECK), "-o", str(TEMPLATE_PPTX)],
        capture_output=True, text=True, encoding="utf-8", errors="replace",
    )
    assert proc.returncode == 2
    assert "template original" in (proc.stderr + proc.stdout)
    # the refusal fired before any write: template bytes unchanged
    assert hashlib.sha256(TEMPLATE_PPTX.read_bytes()).hexdigest() == before


def test_cli_reports_env_errors_as_exit_2(tmp_path):
    # missing explicit manifest -> exit 2, not a traceback / validation code
    missing_manifest = tmp_path / "missing-manifest.json"
    proc = subprocess.run(
        [sys.executable, str(SCRIPT), str(FIXTURE_DECK), "-o", str(tmp_path / "o.pptx"),
         "--manifest", str(missing_manifest)],
        capture_output=True, text=True, encoding="utf-8", errors="replace",
    )
    assert proc.returncode == 2
    assert "Traceback" not in proc.stderr
    # unwritable output directory -> exit 2
    proc = subprocess.run(
        [sys.executable, str(SCRIPT), str(FIXTURE_DECK),
         "-o", str(tmp_path / "no-such-dir" / "o.pptx")],
        capture_output=True, text=True, encoding="utf-8", errors="replace",
    )
    assert proc.returncode == 2
    assert "cannot write output" in proc.stderr


# ---------------------------------------------------------------------------
# T4: the T2 full-archetype fixture (all 6 archetypes, formulas, images,
# captions) renders end to end
# ---------------------------------------------------------------------------

@requires_math
def test_cli_renders_full_fixture(tmp_path):
    out_path, proc = render(tmp_path, FULL_DECK)
    assert proc.returncode == 0, proc.stderr + proc.stdout
    assert out_path.is_file()
    # every archetype in the deck was exercised
    assert "7 page(s)" in proc.stdout


@requires_math
def test_render_full_fixture(tmp_path):
    deck = load("deck.json")
    manifest = json.loads(MANIFEST_PATH.read_text(encoding="utf-8"))
    arch = manifest["archetypes"]
    result = rp.render_deck(deck, manifest, TEMPLATE_PPTX, image_root=FIXTURE_DIR)
    prs = result.presentation

    # -- page count, order, layouts
    assert len(prs.slides) == len(deck["pages"]) == 7
    assert [s.slide_layout.name for s in prs.slides] == [
        "三logo标题页", "标题和内容", "标题和内容",
        "标题和内容", "标题和内容", "标题和内容", "三logo标题页",
    ]

    # -- every page carries its deck title (pages with a title placeholder:
    #    cover, text-formula, text-image, chart-focus, closing; the agenda
    #    page's title block lands in the rebuilt label box, asserted below)
    for i, page in enumerate(deck["pages"]):
        if page["archetype"] == "agenda":
            continue
        title = by_ph(prs.slides[i], "title")
        assert title is not None, f"slide {i + 1} lost its title placeholder"
        assert title.text_frame.text == page["blocks"][0]["text"], f"slide {i + 1} title"

    # -- page 2: agenda = rebuilt safe geometry, no title placeholder at all
    agenda = prs.slides[1]
    assert by_ph(agenda, "title") is None  # slide 2 carries no title placeholder
    assert by_ph(agenda, "sldNum") is not None
    boxes = non_placeholder_shapes(agenda)
    assert len(boxes) == 2
    label, item_list = sorted(boxes, key=lambda sp: sp.top)
    regions = arch["agenda"]["regions"]
    assert_region(label, regions["label"])
    assert label.text_frame.text == deck["pages"][1]["blocks"][0]["text"]
    assert run_style(label) == (40, True, "000000", NOTO, NOTO)
    assert_region(item_list, regions["list"])
    assert item_list.text_frame.text == "\n".join(deck["pages"][1]["blocks"][1]["items"])
    assert run_style(item_list) == (22, False, "000000", NOTO, NOTO)
    # the acceptance check: nothing on the agenda page leaves the 13.33x7.5 canvas
    canvas_w, canvas_h = Inches(13.3333), Inches(7.5)
    for sp in shapes(agenda):
        assert sp.left >= 0 and sp.top >= 0
        assert sp.left + sp.width <= canvas_w + 2
        assert sp.top + sp.height <= canvas_h + 2
    # the safe geometry replaces the original overflow: list bottom at 6.8, not 10.0
    assert item_list.top + item_list.height == pytest.approx(Inches(6.8), abs=2)

    # -- page 3 still carries its native OMML formulas
    assert omath_count(prs.slides[2]) == 2

    # -- page 5: text-image = subhead + text/list + fitted images + caption
    ti = prs.slides[4]
    ti_page = deck["pages"][4]
    ti_regions = arch["text-image"]["regions"]
    title = by_ph(ti, "title")
    assert title.text_frame.text == ti_page["blocks"][0]["text"]
    assert_region(title, ti_regions["title"])
    boxes = non_placeholder_shapes(ti)
    subhead = next(sp for sp in boxes if sp.name == "SubheadArea")
    assert_region(subhead, ti_regions["subhead"])
    assert subhead.text_frame.text == ti_page["blocks"][1]["text"]
    assert run_style(subhead) == (24, True, "548235", NOTO, NOTO)
    textbox = next(sp for sp in boxes if sp.name == "TextArea")
    assert_region(textbox, ti_regions["text"])
    assert textbox.text_frame.text == "\n".join(
        [ti_page["blocks"][2]["text"], *(f"- {item}" for item in ti_page["blocks"][3]["items"])]
    )
    assert run_style(textbox) == (18, False, "595959", NOTO, NOTO)
    pics = [sp for sp in shapes(ti) if sp.shape_type == 13]  # PICTURE
    assert len(pics) == 2
    assert_fitted_in(pics[0], ti_regions["image_slots"][0])
    assert_fitted_in(pics[1], ti_regions["image_slots"][1])
    # colored fixture images carry no hairline (white-background images do)
    for pic in pics:
        assert pic._element.spPr.find(f"{{{NS_A}}}ln") is None
    # caption: 12pt white-on-black-scrim strip over the fitted image's bottom
    # edge (hugs the fitted picture, not the slot)
    caption = next(sp for sp in boxes if sp.name == "Caption1")
    strip_h = Inches(0.30)
    assert list(region_of(caption)) == pytest.approx(
        [pics[0].left, pics[0].top + pics[0].height - strip_h,
         pics[0].width, strip_h], abs=2)
    assert caption.text_frame.text == ti_page["blocks"][4]["caption"]
    assert run_style(caption) == (12, False, "FFFFFF", NOTO, NOTO)
    scrim = caption._element.find(
        f"{{{NS_P}}}spPr/{{{NS_A}}}solidFill/{{{NS_A}}}srgbClr")
    assert scrim is not None and scrim.get("val") == "000000"
    alpha = scrim.find(f"{{{NS_A}}}alpha")
    assert alpha is not None and alpha.get("val") == "60000"
    assert len(boxes) == 5  # subhead + text + caption + 2 pictures

    # -- page 6: chart-focus = one big fitted image + side comment
    cf = prs.slides[5]
    cf_regions = arch["chart-focus"]["regions"]
    assert by_ph(cf, "title").text_frame.text == deck["pages"][5]["blocks"][0]["text"]
    pics = [sp for sp in shapes(cf) if sp.shape_type == 13]
    assert len(pics) == 1
    assert_fitted_in(pics[0], cf_regions["chart"])
    comment = next(sp for sp in non_placeholder_shapes(cf) if sp.name == "CommentArea")
    assert_region(comment, cf_regions["comment"])
    assert comment.text_frame.text == deck["pages"][5]["blocks"][2]["text"]
    assert run_style(comment) == (18, False, "595959", NOTO, NOTO)

    # -- page 7: closing = cover variant
    closing = prs.slides[6]
    title = by_ph(closing, "title")
    assert title.text_frame.text == deck["pages"][6]["blocks"][0]["text"]
    assert_region(title, arch["closing"]["regions"]["title"])
    assert run_style(title) == TITLE_44
    sub = by_ph(closing, "body")
    assert sub.text_frame.text == deck["pages"][6]["blocks"][1]["text"]
    assert_region(sub, arch["closing"]["regions"]["subtitle"])
    assert by_ph(closing, "sldNum") is None

    # -- slide size untouched
    assert (prs.slide_width, prs.slide_height) == (12192000, 6858000)

    # -- zero template residue anywhere in the deck
    for slide in prs.slides:
        text = all_text(slide)
        for residue in TEMPLATE_RESIDUE:
            assert residue not in text

    # -- the deck round-trips through a save
    prs.save(tmp_path / "full.pptx")


@requires_math
def test_full_fixture_template_original_is_never_modified(tmp_path):
    before = hashlib.sha256(TEMPLATE_PPTX.read_bytes()).hexdigest()
    _, proc = render(tmp_path, FULL_DECK)
    assert proc.returncode == 0, proc.stderr + proc.stdout
    assert hashlib.sha256(TEMPLATE_PPTX.read_bytes()).hexdigest() == before


def test_text_image_single_image_uses_primary_column():
    # 单图变体：图片入 image_primary 大区（右栏整版），文字入 text_column（左栏），
    # 不再走模板原 image_slots 的横条小图位
    deck = load("deck.json")
    manifest = json.loads(MANIFEST_PATH.read_text(encoding="utf-8"))
    page = deck["pages"][4]  # text-image page
    keep = page["blocks"][4]  # first image block (keep exactly one)
    page["blocks"] = [page["blocks"][0], page["blocks"][1], page["blocks"][2], keep]
    result = rp.render_deck(deck, manifest, TEMPLATE_PPTX, image_root=FIXTURE_DIR)
    slide = result.presentation.slides[4]
    regions = manifest["archetypes"]["text-image"]["regions"]
    pics = [sp for sp in shapes(slide) if sp.shape_type == 13]
    assert len(pics) == 1
    assert_fitted_in(pics[0], regions["image_primary"])
    column = next(sp for sp in non_placeholder_shapes(slide) if sp.name == "TextColumn")
    assert_region(column, regions["text_column"])
    assert not any(sp.name == "TextArea" for sp in non_placeholder_shapes(slide))


# ---------------------------------------------------------------------------
# S3: typography.tokens land at run level — emphasis runs, spacing rhythm,
# token insets, caption scrim, white-image hairline, font embed
# ---------------------------------------------------------------------------

@requires_math
def test_emphasis_marker_renders_green_bold_inline_runs():
    deck = load("deck_pptx.json")
    manifest = json.loads(MANIFEST_PATH.read_text(encoding="utf-8"))
    result = rp.render_deck(deck, manifest, TEMPLATE_PPTX)
    box = next(sp for sp in non_placeholder_shapes(result.presentation.slides[2])
               if sp.name == "ContentArea")
    runs = shape_runs(box)
    # markers never survive into a:t text; the keyword becomes its own run
    assert all("**" not in (r[0] or "") for r in runs)
    assert ("对比学习", 20, True, "548235") in runs
    tail = next(r for r in runs if r[0].startswith("通过拉近正例"))
    assert tail[1:] == (20, False, "000000")
    # the marker-free paragraph stays one regular black run
    assert ("神经机器翻译在低资源场景下性能受限，平行语料稀缺是主要瓶颈。",
            20, False, "000000") in runs


def test_unmatched_emphasis_marker_stays_literal():
    deck = load("deck_pptx.json")
    deck["pages"][2]["blocks"] = [
        {"type": "title", "text": "研究背景与动机"},
        {"type": "text", "text": "不带成对标记的 ** 星号正文"},   # 1 marker
        {"type": "text", "text": "**** 四个星号"},               # empty pair
    ]
    manifest = json.loads(MANIFEST_PATH.read_text(encoding="utf-8"))
    result = rp.render_deck(deck, manifest, TEMPLATE_PPTX)
    box = next(sp for sp in non_placeholder_shapes(result.presentation.slides[2])
               if sp.name == "ContentArea")
    # unpaired and empty markers both stay literal — never an empty paragraph
    assert [(r[0], r[2]) for r in shape_runs(box)] == [
        ("不带成对标记的 ** 星号正文", False),
        ("**** 四个星号", False),
    ]


@requires_math
def test_emphasis_is_body_only_and_covers_body_flow_roles():
    # 正文内 convention: ** in a title/subtitle stays literal black (no green
    # run); the same marker in body-flow text — body paragraphs, agenda list
    # items — emphasizes (explicitly accepted extension of 正文: every role
    # that carries the rhythm)
    deck = load("deck.json")
    deck["pages"][0]["blocks"][0]["text"] = "带**标记**的封面标题"
    deck["pages"][0]["blocks"][1]["text"] = "汇报人:张三 **强调** 2026-08-25"
    deck["pages"][3]["blocks"][2]["text"] = "**关键词**正文一行。"
    deck["pages"][1]["blocks"][1]["items"] = ["一 **重点**章节", "二 次要章节"]
    # literal roles: subhead and caption keep ** verbatim (the caption case
    # is what raw font-subset collection depends on)
    deck["pages"][4]["blocks"][1]["text"] = "小标题**不**强调"
    deck["pages"][4]["blocks"][4]["caption"] = "图注**保持**字面"
    manifest = json.loads(MANIFEST_PATH.read_text(encoding="utf-8"))
    result = rp.render_deck(deck, manifest, TEMPLATE_PPTX, image_root=FIXTURE_DIR)
    prs = result.presentation

    cover_title = by_ph(prs.slides[0], "title")
    assert shape_runs(cover_title) == [("带**标记**的封面标题", 44, True, "000000")]
    cover_sub = by_ph(prs.slides[0], "body")
    assert shape_runs(cover_sub) == [("汇报人:张三 **强调** 2026-08-25",
                                      24, False, "000000")]
    body_box = next(sp for sp in non_placeholder_shapes(prs.slides[3])
                    if sp.name == "ContentArea")
    assert ("关键词", 20, True, "548235") in shape_runs(body_box)
    agenda_list = next(sp for sp in non_placeholder_shapes(prs.slides[1])
                       if sp.name == "AgendaList")
    assert ("重点", 22, True, "548235") in shape_runs(agenda_list)
    ti_boxes = non_placeholder_shapes(prs.slides[4])
    subhead = next(sp for sp in ti_boxes if sp.name == "SubheadArea")
    assert shape_runs(subhead) == [("小标题**不**强调", 24, True, "548235")]
    caption = next(sp for sp in ti_boxes if sp.name == "Caption1")
    assert shape_runs(caption) == [("图注**保持**字面", 12, False, "FFFFFF")]


@requires_math
def test_body_rhythm_and_insets_follow_tokens():
    deck = load("deck.json")
    manifest = json.loads(MANIFEST_PATH.read_text(encoding="utf-8"))
    result = rp.render_deck(deck, manifest, TEMPLATE_PPTX, image_root=FIXTURE_DIR)
    prs = result.presentation

    def by_name(slide, name):
        return next(sp for sp in non_placeholder_shapes(slide) if sp.name == name)

    # flowing body text carries the rhythm: line spacing 1.5, before 12pt
    # (formula-first flows open with the OMML shell's 130%/16pt rhythm)
    flow2 = by_name(prs.slides[2], "ContentArea")
    assert first_para_props(flow2) == (130, 16)
    assert prose_para_props(flow2) == (150, 12)
    assert first_para_props(by_name(prs.slides[3], "ContentArea")) == (150, 12)
    assert first_para_props(by_name(prs.slides[1], "AgendaList")) == (150, 12)
    assert first_para_props(by_name(prs.slides[4], "TextArea")) == (150, 12)
    assert first_para_props(by_name(prs.slides[5], "CommentArea")) == (150, 12)
    # ceremonial single-line roles stay at single spacing, no space-before
    assert first_para_props(by_ph(prs.slides[0], "title")) == (100, 0)
    assert first_para_props(by_ph(prs.slides[0], "body")) == (100, 0)
    assert first_para_props(by_name(prs.slides[1], "AgendaLabel")) == (100, 0)
    assert first_para_props(by_name(prs.slides[4], "SubheadArea")) == (100, 0)
    assert first_para_props(by_name(prs.slides[4], "Caption1")) == (100, 0)
    # synthesized content boxes carry the token insets (0.1in sides, 0.05in ends)
    for slide in prs.slides:
        for sp in non_placeholder_shapes(slide):
            if sp.shape_type == 13:  # pictures carry no bodyPr
                continue
            assert body_insets(sp) == (91440, 45720, 91440, 45720), sp.name


@requires_math
def test_hairline_outlines_white_background_images():
    deck = load("deck.json")
    deck["pages"][5]["blocks"][1]["path"] = "material/images/white-chart.png"
    manifest = json.loads(MANIFEST_PATH.read_text(encoding="utf-8"))
    result = rp.render_deck(deck, manifest, TEMPLATE_PPTX, image_root=FIXTURE_DIR)
    pic = next(sp for sp in shapes(result.presentation.slides[5])
               if sp.shape_type == 13)
    ln = pic._element.spPr.find(f"{{{NS_A}}}ln")
    assert ln is not None and ln.get("w") == "9525"  # 0.75pt
    clr = ln.find(f"{{{NS_A}}}solidFill/{{{NS_A}}}srgbClr")
    assert clr is not None and clr.get("val") == "D9D9D9"


@requires_math
def test_full_fixture_carries_embedded_font_subsets(tmp_path):
    out_path, proc = render(tmp_path, FULL_DECK)
    assert proc.returncode == 0, proc.stderr + proc.stdout
    assert "embedded Noto Sans SC subset" in proc.stdout
    with zipfile.ZipFile(out_path) as z:
        names = z.namelist()
        assert [n for n in names if n.endswith(".fntdata")] == [
            "ppt/fonts/font1.fntdata", "ppt/fonts/font2.fntdata"]
        assert b'<p:font typeface="Noto Sans SC"' in z.read("ppt/presentation.xml")
